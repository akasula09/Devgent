"""
Devgent — Web Backend (Vercel Python Serverless Function)
===========================================================
Deployed at:  /api/AIProcess  (Vercel auto-routes api/AIProcess.py -> /api/AIProcess)
Called by:    chat.html

This is the browser-facing counterpart to the Devgent desktop app. Same job —
turn a plain-English request into a real Word doc / PowerPoint / spreadsheet /
PDF / code file — but re-architected for a stateless, public serverless
function instead of a trusted desktop process:

  * No `run_command` / `start_dev_server` actions. Those meant "execute an
    arbitrary shell command" on the desktop, where the *user* is the one
    running it on their own machine. Exposing that over a public HTTP
    endpoint would let anyone who finds this URL run arbitrary code on
    Vercel's infrastructure — a critical remote-code-execution hole. So the
    web version only ever *generates files* and hands them back as
    downloads; it never executes anything.
  * No server-side memory store. Vercel functions don't persist state
    between invocations, so "memory" is round-tripped with the client on
    every request instead of written to disk (see the `memory` field below).
  * The Gemini API key is supplied by the browser on every request and is
    never logged, stored, or written anywhere server-side.

Request body (POST, JSON):
{
  "message":  "user's chat message",
  "history":  [{"role": "user"|"model", "content": "..."}, ...]   # prior turns, most recent last, EXCLUDING this message
  "apiKey":   "user's Gemini API key",
  "model":    "gemini-3.6-flash" | "gemini-3.5-flash-lite",
  "memory":   { ...arbitrary persisted key/values... },
  "attachment": {"filename": "notes.docx", "data_base64": "..."}   # optional
}

Response body (JSON):
{
  "reply":  "assistant's chat text (markdown, JSON action block stripped)",
  "files":  [{"filename": "...", "mimetype": "...", "data_base64": "..."}],
  "memory": { ...updated memory dict, echo this back on the next request... }
}
"""

from http.server import BaseHTTPRequestHandler
import json
import os
import re
import io
import zipfile
import base64
import traceback
from datetime import date

import requests

GEMINI_API_BASE = "https://generativelanguage.googleapis.com/v1beta/models"

GEMINI_API_KEY_HELP_TEXT = """Getting a Gemini API key is completely free and only takes a minute:

1. Open Google AI Studio in your browser at https://aistudio.google.com/apikey
2. Sign in with your Google account if you're asked to.
3. Click "Create API key", then choose to create a new project (or pick an existing one) when prompted.
4. Copy the API key that gets generated.
5. Paste it into the API key field at the top of the page.

No credit card or payment is required — the free tier is enough to chat with Devgent."""

SYSTEM_INSTRUCTION_TEMPLATE = """
You are Devgent, an AI assistant running as a web app in the person's browser.
Current date: {current_date}

=== CORE OPERATING PRINCIPLES ===
1. Respond clearly to the user in plain chat text, then, ONLY if you are actually taking an action this turn, attach a single JSON block at the very end of your response inside triple backticks tagged with json.
2. The JSON block can contain an "actions" list so you can queue multiple actions in a single turn.
3. File Naming: Always specify exact file extensions (.py, .js, .html, .css, .json, .txt, .pdf, .docx, .xlsx, .pptx).
4. Environment Awareness: You are running as a stateless web backend. You have NO terminal access and CANNOT run commands, install packages, or start servers. You can only generate files, which are handed to the user as downloads in their browser. Never claim to run or execute code — only write it.
5. No Permission Theater Needed: Unlike a desktop app, generating a file here is harmless (nothing on the user's machine is touched until they click Download), so you do not need to ask "can I create this file?" before emitting a create_* action — just briefly describe what you're making, then issue the action.
6. IMPORTANT — Never claim a file has been "saved" or "created on your computer": this is a stateless web backend with no access to the user's file system. When a create_* action finishes, a download card with a Download button appears automatically in the chat — you do not need to mention how to download it, but you must phrase your own reply as "I've put together..." / "Here's your..." / "Ready below —", never "I've saved this to your Documents folder" or any wording implying direct disk access.
7. Markdown In Chat: Your plain-text reply (everything outside the JSON block) is rendered as real Markdown — headings, **bold**, *italics*, bullet/numbered lists, tables, and fenced code blocks with a language tag (e.g. ```python ... ```) will all display properly, including syntax-highlighted code blocks. This is completely separate from the JSON action block used to trigger file creation.

=== CLARIFY BEFORE ACTING (VERY IMPORTANT) ===
If the user's request is vague, incomplete, or does not tell you concretely WHAT to build, write, or do — for example "I want to make a python script", "build me something cool", "make a presentation", "help me with a document" — then you must NOT take any action and must NOT emit a JSON action block at all this turn. Instead:
- Ask the user, in plain chat text only, exactly what they want: purpose, subject matter, desired content/sections, style, filename, etc. — whatever is missing.
- If they already gave you some partial detail, save it to memory with "manage_memory" so it isn't lost (this is the one exception where you may still emit a JSON block even without creating a file), e.g.:
```json
{{
  "actions": [
    {{"action": "manage_memory", "action_data": {{"updates": {{"pending_request": "User wants a Python script, no further details yet."}}}}}}
  ]
}}
```
- Wait for the user's next message before creating anything.
Only move on to create_txt / create_docx / create_pdf / create_pptx / create_xlsx / create_batch_files once the user has given you enough concrete specifics to act on.

=== DYNAMIC MEMORY STORE (BE AGGRESSIVE ABOUT USING THIS) ===
Current Stored Memories:
{memory_state}

You can save, update, or discard memory items dynamically to retain user context across conversations. Treat this as a first-class habit — on almost every turn, ask yourself "did I just learn something worth remembering?" That includes, but is not limited to:
- The user's name, role, or how they prefer to be addressed.
- Project name, purpose, and key details already established.
- Stylistic or technical preferences (tone, color themes, frameworks).
- Decisions the user has made or corrections they've given you (these should overwrite/replace old memory values, not just add new ones).
- The current status of a multi-step task, so you can resume correctly next session.
If the answer is yes, include a "manage_memory" action alongside any other action in the same "actions" array. Err on the side of saving too much rather than too little.

=== MULTI-ACTION SCHEMA ===
```json
{{
  "actions": [
    {{"action": "create_txt", "action_data": {{"filename": "app.py", "content": "print('hello')"}}}}
  ]
}}
```

=== AVAILABLE ACTIONS ===

1. Manage Dynamic Memory (manage_memory)
{{"action": "manage_memory", "action_data": {{"updates": {{"user_theme": "dark"}}, "discards": ["old_project_key"]}}}}

2. Advanced Styled PowerPoint Presentation (create_pptx)
```json
{{
  "action": "create_pptx",
  "action_data": {{
    "filename": "pitch.pptx",
    "title": "Deck Title",
    "subtitle": "Optional subtitle line",
    "title_color": "#3B82F6",
    "subtitle_color": "#94A3B8",
    "background_color": "#0F172A",
    "theme_accent": "#3B82F6",
    "slides": [
      {{
        "layout": "bullets",
        "title": "Slide Heading",
        "background_color": "#F8FAFC",
        "header_color": "#1E293B",
        "header_text_color": "#FFFFFF",
        "accent_color": "#3B82F6",
        "content": ["First point", "Second point", "Third point"],
        "chart": {{
          "type": "column",
          "categories": ["Q1", "Q2", "Q3"],
          "series": [{{"name": "Revenue", "values": [10, 20, 30]}}]
        }}
      }},
      {{
        "layout": "two_column",
        "title": "Comparison",
        "left_heading": "Pros",
        "content_left": ["Fast", "Cheap"],
        "right_heading": "Cons",
        "content_right": ["Riskier"]
      }},
      {{
        "layout": "quote",
        "title": "Customer Voice",
        "quote_text": "This product changed how we work.",
        "quote_author": "Jane Doe, CTO"
      }},
      {{
        "layout": "section",
        "title": "Part Two: Results",
        "subtitle": "Optional divider subtitle",
        "background_color": "#3B82F6",
        "header_text_color": "#FFFFFF"
      }}
    ]
  }}
}}
```
Rules:
- "layout" is one of: "bullets" (default; supports an optional "chart" side-by-side), "two_column", "quote", or "section" (a full-bleed divider slide).
- All colors are hex strings like "#0F172A". Any slide can override background_color, header_color, header_text_color, and accent_color; otherwise it inherits "theme_accent".
- Chart "type" is one of "bar", "column", "pie", "line".
- Prefer 3-6 short bullets per slide rather than long paragraphs.

3. Excel Spreadsheet (create_xlsx)
{{"action": "create_xlsx", "action_data": {{"filename": "data.xlsx", "sheet_name": "Metrics", "headers": ["Name", "Score"], "rows": [["Alice", 95]]}}}}

4. Word Document (create_docx)
{{"action": "create_docx", "action_data": {{"filename": "doc.docx", "title": "Proposal", "heading": "Summary", "paragraphs": ["Details..."]}}}}

5. Styled PDF Document (create_pdf)
```json
{{
  "action": "create_pdf",
  "action_data": {{
    "filename": "report.pdf",
    "title": "Report Title",
    "subtitle": "Optional subtitle",
    "accent_color": "#3B82F6",
    "sections": [
      {{
        "heading": "Overview",
        "paragraphs": ["Body paragraph text goes here."],
        "bullets": ["Point one", "Point two"],
        "quote": "An optional pull-quote highlighting something important.",
        "table": {{"headers": ["Name", "Score"], "rows": [["Alice", "95"], ["Bob", "88"]]}}
      }}
    ]
  }}
}}
```
- For very simple documents you may skip "sections" and just pass a flat "paragraphs" list.

6. Single Text/Code File (create_txt)
{{"action": "create_txt", "action_data": {{"filename": "app.py", "content": "print('hello')"}}}}

7. Batch Files, packaged as one .zip (create_batch_files)
{{"action": "create_batch_files", "action_data": {{"zip_filename": "site.zip", "files": [{{"filename": "index.html", "content": "<h1>Hello</h1>"}}]}}}}

Remember to USE your memory as much as you can. You will remember absolutely nothing that you do not store in your memory.
"""


# =============================================================================
# Small color helpers (shared by docx/pptx/pdf/xlsx builders)
# =============================================================================

def hex_to_rgb_tuple(hex_str):
    hex_str = (hex_str or "").lstrip('#')
    if len(hex_str) != 6:
        hex_str = "3B82F6"
    return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)


def shade_color(hex_str, amount):
    r, g, b = hex_to_rgb_tuple(hex_str)
    if amount >= 0:
        r = int(r + (255 - r) * amount)
        g = int(g + (255 - g) * amount)
        b = int(b + (255 - b) * amount)
    else:
        r, g, b = int(r * (1 + amount)), int(g * (1 + amount)), int(b * (1 + amount))
    r, g, b = max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b))
    return f"#{r:02X}{g:02X}{b:02X}"


def safe_filename(raw, default_ext):
    name = os.path.basename((raw or "").strip() or f"file{default_ext}")
    name = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "_", name)
    if not name.lower().endswith(default_ext):
        name += default_ext
    return name


# =============================================================================
# File builders — each returns raw bytes
# =============================================================================

def build_docx_bytes(data):
    import docx
    from docx.shared import Pt, RGBColor, Inches, Emu
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    doc = docx.Document()
    title_text = data.get("title", "Document Title")
    p_title = doc.add_paragraph()
    run = p_title.add_run(title_text)
    run.font.name = 'Calibri'
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

    accent_table = doc.add_table(rows=1, cols=1)
    accent_cell = accent_table.rows[0].cells[0]
    accent_cell.width = Inches(6.0)
    shading = OxmlElement('w:shd')
    shading.set(qn('w:fill'), "3B82F6")
    accent_cell._tc.get_or_add_tcPr().append(shading)
    accent_table.rows[0].height = Emu(9525 * 3)

    if data.get("heading"):
        p_head = doc.add_paragraph()
        r_head = p_head.add_run(data["heading"])
        r_head.font.name = 'Calibri'
        r_head.font.size = Pt(16)
        r_head.font.bold = True
        r_head.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)

    for ptext in data.get("paragraphs", []):
        p = doc.add_paragraph()
        r = p.add_run(ptext)
        r.font.name = 'Calibri'
        r.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pdf_bytes(data):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    )

    accent_hex = data.get("accent_color", "#3B82F6")
    accent = colors.HexColor(accent_hex)
    accent_dark = colors.HexColor(shade_color(accent_hex, -0.35))
    text_dark = colors.HexColor("#1F2937")
    muted = colors.HexColor("#64748B")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter,
        topMargin=0.9 * inch, bottomMargin=0.8 * inch,
        leftMargin=0.75 * inch, rightMargin=0.75 * inch
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('T', parent=styles['Heading1'], fontSize=26, leading=30, textColor=colors.white, spaceAfter=4, fontName="Helvetica-Bold")
    subtitle_style = ParagraphStyle('S', parent=styles['Normal'], fontSize=12, leading=16, textColor=colors.HexColor("#E2E8F0"))
    heading_style = ParagraphStyle('H', parent=styles['Heading2'], fontSize=15, leading=19, textColor=accent_dark, spaceBefore=18, spaceAfter=8, fontName="Helvetica-Bold")
    body_style = ParagraphStyle('B', parent=styles['Normal'], fontSize=10.5, leading=15, textColor=text_dark, spaceAfter=8)
    bullet_style = ParagraphStyle('BL', parent=body_style, leftIndent=14, spaceAfter=6)
    quote_style = ParagraphStyle('Q', parent=body_style, fontName="Helvetica-Oblique", leftIndent=16, textColor=muted)

    story = []
    title_text = data.get("title", "Generated Report")
    subtitle_text = data.get("subtitle", "")
    title_cell = [Paragraph(title_text, title_style)]
    if subtitle_text:
        title_cell.append(Paragraph(subtitle_text, subtitle_style))
    title_table = Table([[title_cell]], colWidths=[6.9 * inch])
    title_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), accent),
        ('LEFTPADDING', (0, 0), (-1, -1), 18), ('RIGHTPADDING', (0, 0), (-1, -1), 18),
        ('TOPPADDING', (0, 0), (-1, -1), 18), ('BOTTOMPADDING', (0, 0), (-1, -1), 18),
    ]))
    story.append(title_table)
    story.append(Spacer(1, 18))

    sections = data.get("sections") or [{"paragraphs": data.get("paragraphs", [])}]
    for section in sections:
        if section.get("heading"):
            h_table = Table([[Paragraph(section["heading"], heading_style)]], colWidths=[6.9 * inch])
            h_table.setStyle(TableStyle([
                ('LINEBELOW', (0, 0), (-1, -1), 1.4, accent),
                ('LEFTPADDING', (0, 0), (-1, -1), 0), ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            story.append(h_table)

        for paragraph in section.get("paragraphs", []):
            story.append(Paragraph(paragraph, body_style))

        bullets = section.get("bullets", [])
        if bullets:
            items = [ListItem(Paragraph(b, bullet_style), bulletColor=accent, value='●') for b in bullets]
            story.append(ListFlowable(items, bulletType='bullet', start='●', leftIndent=16))
            story.append(Spacer(1, 6))

        if section.get("quote"):
            q_table = Table([[Paragraph(f"\u201c{section['quote']}\u201d", quote_style)]], colWidths=[6.9 * inch])
            q_table.setStyle(TableStyle([
                ('LINEBEFORE', (0, 0), (0, -1), 3, accent),
                ('LEFTPADDING', (0, 0), (-1, -1), 14),
                ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ]))
            story.append(q_table)
            story.append(Spacer(1, 8))

        table_data = section.get("table")
        if table_data and table_data.get("headers") and table_data.get("rows"):
            tbl = Table([table_data["headers"]] + table_data["rows"], hAlign='LEFT')
            tbl.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), accent),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9.5),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F5F9")]),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 10))
        story.append(Spacer(1, 4))

    def _footer(canvas, doc_):
        canvas.saveState()
        canvas.setStrokeColor(accent)
        canvas.setLineWidth(1)
        canvas.line(0.75 * inch, 0.6 * inch, letter[0] - 0.75 * inch, 0.6 * inch)
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(muted)
        canvas.drawString(0.75 * inch, 0.45 * inch, title_text)
        canvas.drawRightString(letter[0] - 0.75 * inch, 0.45 * inch, f"Page {doc_.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=_footer, onLaterPages=_footer)
    return buf.getvalue()


def build_xlsx_bytes(data):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = data.get("sheet_name", "Sheet1")

    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    alt_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
    thin_border = Border(*(Side(style='thin', color='CBD5E1') for _ in range(4)))

    headers = data.get("headers", [])
    if headers:
        ws.append(headers)
        for col_num in range(1, len(headers) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

    rows = data.get("rows", [])
    for row_idx, row in enumerate(rows, start=2):
        ws.append(row)
        for col_num in range(1, len(row) + 1):
            cell = ws.cell(row=row_idx, column=col_num)
            cell.font = Font(name="Calibri", size=10)
            cell.border = thin_border
            if row_idx % 2 == 0:
                cell.fill = alt_fill

    for col in ws.columns:
        max_len = max((len(str(c.value or '')) for c in col), default=0)
        ws.column_dimensions[get_column_letter(col[0].column)].width = max(max_len + 4, 12)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _add_pptx_footer(slide, index, total, accent_color, slide_w, slide_h, on_light_bg=True):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), slide_h - Inches(0.45), slide_w - Inches(1.0), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(accent_color)
    line.line.fill.background()
    line.shadow.inherit = False

    tb = slide.shapes.add_textbox(slide_w - Inches(1.6), slide_h - Inches(0.42), Inches(1.1), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{index:02d} / {total:02d}"
    p.font.size = Pt(10)
    p.font.color.rgb = _rgb(accent_color if on_light_bg else "#94A3B8")
    p.alignment = PP_ALIGN.RIGHT


def _rgb(hex_str):
    from pptx.dml.color import RGBColor
    r, g, b = hex_to_rgb_tuple(hex_str)
    return RGBColor(r, g, b)


def build_pptx_bytes(data):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
    theme_accent = data.get("theme_accent", data.get("title_color", "#3B82F6"))

    # ---- Title slide ----
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bg_hex = data.get("background_color", "#0F172A")
    bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = _rgb(title_bg_hex); bg.line.fill.background(); bg.shadow.inherit = False

    accent_circle = title_slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(4.0), Inches(-2.0), Inches(6.0), Inches(6.0))
    accent_circle.fill.solid(); accent_circle.fill.fore_color.rgb = _rgb(shade_color(theme_accent, -0.35))
    accent_circle.line.fill.background(); accent_circle.shadow.inherit = False

    accent_bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), Inches(0.18), Inches(1.9))
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = _rgb(theme_accent)
    accent_bar.line.fill.background(); accent_bar.shadow.inherit = False

    txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(9.5), Inches(3.0))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Presentation Title")
    p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = _rgb(data.get("title_color", theme_accent))

    sub = tf.add_paragraph()
    sub.text = data.get("subtitle", "Generated by Devgent")
    sub.font.size = Pt(20)
    sub.font.color.rgb = _rgb(data.get("subtitle_color", "#94A3B8"))
    sub.space_before = Pt(14)

    slides_list = data.get("slides", [])
    total_slides = len(slides_list)

    for i, sdata in enumerate(slides_list, start=1):
        layout = str(sdata.get("layout", "bullets")).lower()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        accent_color = sdata.get("accent_color", theme_accent)
        slide_bg_hex = sdata.get("background_color", "#3B82F6" if layout == "section" else "#F8FAFC")

        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        s_bg.fill.solid(); s_bg.fill.fore_color.rgb = _rgb(slide_bg_hex)
        s_bg.line.fill.background(); s_bg.shadow.inherit = False

        is_dark_bg = slide_bg_hex.upper() not in ["#FFFFFF", "#F8FAFC", "#F1F5F9", "#FFF"]
        card_fill = "#1E293B" if is_dark_bg else "#FFFFFF"
        default_text = sdata.get("text_color", "#F8FAFC" if is_dark_bg else "#334155")

        if layout == "section":
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.6))
            tf2 = tb.text_frame; tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sdata.get("title", "Section")
            p2.font.size = Pt(42); p2.font.bold = True
            p2.font.color.rgb = _rgb(sdata.get("header_text_color", "#FFFFFF"))
            p2.alignment = PP_ALIGN.CENTER
            if sdata.get("subtitle"):
                sp = tf2.add_paragraph()
                sp.text = sdata["subtitle"]
                sp.font.size = Pt(18)
                sp.font.color.rgb = _rgb(shade_color(sdata.get("header_text_color", "#FFFFFF"), -0.2))
                sp.alignment = PP_ALIGN.CENTER
            _add_pptx_footer(slide, i, total_slides, "#FFFFFF", SLIDE_W, SLIDE_H, on_light_bg=False)
            continue

        header_bg_hex = sdata.get("header_color", accent_color)
        header_txt_hex = sdata.get("header_text_color", "#FFFFFF")
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), SLIDE_W - Inches(1.0), Inches(0.9))
        banner.fill.solid(); banner.fill.fore_color.rgb = _rgb(header_bg_hex)
        banner.line.fill.background(); banner.shadow.inherit = False
        bp = banner.text_frame.paragraphs[0]
        bp.text = sdata.get("title", "Slide Title")
        bp.font.size = Pt(24); bp.font.bold = True
        bp.font.color.rgb = _rgb(header_txt_hex)

        if layout == "quote":
            quote_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.8), SLIDE_W - Inches(2.8), Inches(4.6))
            quote_card.fill.solid(); quote_card.fill.fore_color.rgb = _rgb(card_fill)
            quote_card.line.fill.background(); quote_card.shadow.inherit = False

            q_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(1.8), Inches(0.12), Inches(4.6))
            q_bar.fill.solid(); q_bar.fill.fore_color.rgb = _rgb(accent_color)
            q_bar.line.fill.background(); q_bar.shadow.inherit = False

            qtf = quote_card.text_frame; qtf.word_wrap = True
            qtf.margin_left = Inches(0.5); qtf.margin_right = Inches(0.4)
            qp = qtf.paragraphs[0]
            qp.text = f"\u201c{sdata.get('quote_text', '')}\u201d"
            qp.font.size = Pt(24); qp.font.italic = True
            qp.font.color.rgb = _rgb(default_text)
            if sdata.get("quote_author"):
                ap = qtf.add_paragraph()
                ap.text = f"\u2014 {sdata['quote_author']}"
                ap.font.size = Pt(16); ap.font.bold = True
                ap.font.color.rgb = _rgb(accent_color)
                ap.space_before = Pt(16)

        elif layout == "two_column":
            col_w = int((SLIDE_W - Inches(1.4)) / 2)
            columns = [
                (Inches(0.5), sdata.get("left_heading", ""), sdata.get("content_left", [])),
                (Inches(0.5) + col_w + Inches(0.4), sdata.get("right_heading", ""), sdata.get("content_right", [])),
            ]
            for left, heading, items in columns:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), col_w, Inches(5.4))
                card.fill.solid(); card.fill.fore_color.rgb = _rgb(card_fill)
                card.line.color.rgb = _rgb(shade_color(card_fill, -0.05))
                card.shadow.inherit = False
                ctf = card.text_frame; ctf.word_wrap = True
                ctf.margin_left = Inches(0.3); ctf.margin_top = Inches(0.25)
                first = False
                if heading:
                    hp = ctf.paragraphs[0]
                    hp.text = heading; hp.font.bold = True; hp.font.size = Pt(16)
                    hp.font.color.rgb = _rgb(accent_color); hp.space_after = Pt(10)
                    first = True
                for bullet in items:
                    bp2 = ctf.add_paragraph() if first else ctf.paragraphs[0]
                    first = True
                    r1 = bp2.add_run(); r1.text = "\u25cf "; r1.font.color.rgb = _rgb(accent_color); r1.font.size = Pt(14)
                    r2 = bp2.add_run(); r2.text = bullet; r2.font.color.rgb = _rgb(default_text); r2.font.size = Pt(14)
                    bp2.space_after = Pt(8)

        else:  # bullets (default), optionally with a chart
            content_items = sdata.get("content", [])
            has_chart = isinstance(sdata.get("chart"), dict)
            text_width = Inches(5.6) if has_chart else SLIDE_W - Inches(1.0)
            chart_left = Inches(6.5)
            chart_width = SLIDE_W - chart_left - Inches(0.5)

            if content_items:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), text_width, Inches(5.4))
                card.fill.solid(); card.fill.fore_color.rgb = _rgb(card_fill)
                card.line.color.rgb = _rgb(shade_color(card_fill, -0.05))
                card.shadow.inherit = False
                ctf = card.text_frame; ctf.word_wrap = True
                ctf.margin_left = Inches(0.35); ctf.margin_top = Inches(0.3)
                for idx, bullet in enumerate(content_items):
                    cp = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
                    r1 = cp.add_run(); r1.text = "\u25cf "; r1.font.size = Pt(15); r1.font.color.rgb = _rgb(accent_color)
                    r2 = cp.add_run(); r2.text = bullet; r2.font.size = Pt(15); r2.font.color.rgb = _rgb(default_text)
                    cp.space_after = Pt(10)

            if has_chart:
                chart_info = sdata["chart"]
                chart_data = CategoryChartData()
                chart_data.categories = chart_info.get("categories", [])
                for s in chart_info.get("series", []):
                    chart_data.add_series(s.get("name", "Series"), tuple(s.get("values", [])))
                chart_type_map = {
                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                    "pie": XL_CHART_TYPE.PIE, "line": XL_CHART_TYPE.LINE
                }
                xl_type = chart_type_map.get(str(chart_info.get("type", "column")).lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
                chart_shape = slide.shapes.add_chart(xl_type, chart_left, Inches(1.5), chart_width, Inches(5.4), chart_data)
                chart = chart_shape.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

        _add_pptx_footer(slide, i, total_slides, accent_color, SLIDE_W, SLIDE_H, on_light_bg=not is_dark_bg)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_batch_zip_bytes(files):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            fname = safe_filename(f.get("filename", "file.txt"), "")
            zf.writestr(fname, f.get("content", ""))
    return buf.getvalue()


# =============================================================================
# Attachment text extraction (mirrors the desktop app's extract_file_content)
# =============================================================================

def extract_file_content(filename, raw_bytes):
    ext = os.path.splitext(filename or "")[1].lower()
    try:
        if ext == '.docx':
            import docx
            doc = docx.Document(io.BytesIO(raw_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    if any(c is not None for c in row):
                        content.append("\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(content)
        elif ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        else:
            return raw_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[Error extracting file content from {filename}: {e}]"


# =============================================================================
# Gemini call (plain REST — no SDK, keeps the deployed function small)
# =============================================================================

def call_gemini(api_key, model, system_instruction, history, user_message):
    contents = []
    for h in history:
        role = "user" if h.get("role") == "user" else "model"
        contents.append({"role": role, "parts": [{"text": h.get("content", "")}]})
    contents.append({"role": "user", "parts": [{"text": user_message}]})

    body = {
        "contents": contents,
        "systemInstruction": {"parts": [{"text": system_instruction}]},
        "generationConfig": {"temperature": 0.3},
    }
    url = f"{GEMINI_API_BASE}/{model}:generateContent?key={api_key}"
    resp = requests.post(url, json=body, timeout=55)
    if resp.status_code != 200:
        detail = ""
        try:
            detail = resp.json().get("error", {}).get("message", "")
        except Exception:
            detail = resp.text[:300]
        raise RuntimeError(f"Gemini API error ({resp.status_code}): {detail}")

    data = resp.json()
    candidates = data.get("candidates") or []
    if not candidates:
        raise RuntimeError("Gemini returned no candidates (the response may have been blocked by safety filters).")
    parts = candidates[0].get("content", {}).get("parts", [])
    return "".join(p.get("text", "") for p in parts)


# =============================================================================
# Action execution — turns the model's JSON block into downloadable files
# =============================================================================

def execute_actions(actions, memory):
    files = []
    system_notes = []

    for item in actions:
        action_type = item.get("action")
        action_data = item.get("action_data", {}) or {}

        try:
            if action_type == "manage_memory":
                for k in action_data.get("discards", []):
                    memory.pop(k, None)
                for k, v in (action_data.get("updates", {}) or {}).items():
                    memory[k] = v

            elif action_type == "create_txt":
                fname = safe_filename(action_data.get("filename", "output.txt"), "")
                files.append({"filename": fname, "mimetype": "text/plain",
                              "data_base64": base64.b64encode(action_data.get("content", "").encode("utf-8")).decode()})

            elif action_type == "create_docx":
                fname = safe_filename(action_data.get("filename", "Document.docx"), ".docx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                              "data_base64": base64.b64encode(build_docx_bytes(action_data)).decode()})

            elif action_type == "create_pdf":
                fname = safe_filename(action_data.get("filename", "Document.pdf"), ".pdf")
                files.append({"filename": fname, "mimetype": "application/pdf",
                              "data_base64": base64.b64encode(build_pdf_bytes(action_data)).decode()})

            elif action_type == "create_xlsx":
                fname = safe_filename(action_data.get("filename", "Spreadsheet.xlsx"), ".xlsx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                              "data_base64": base64.b64encode(build_xlsx_bytes(action_data)).decode()})

            elif action_type == "create_pptx":
                fname = safe_filename(action_data.get("filename", "Presentation.pptx"), ".pptx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                              "data_base64": base64.b64encode(build_pptx_bytes(action_data)).decode()})

            elif action_type == "create_batch_files":
                zip_name = safe_filename(action_data.get("zip_filename", "devgent_files.zip"), ".zip")
                files.append({"filename": zip_name, "mimetype": "application/zip",
                              "data_base64": base64.b64encode(build_batch_zip_bytes(action_data.get("files", []))).decode()})

            elif action_type in ("run_command", "start_dev_server"):
                system_notes.append(
                    f"⚠️ The `{action_type}` action isn't available in the web version for security reasons "
                    f"(a public server can't safely run arbitrary commands). Download the Devgent desktop app "
                    f"for terminal and dev-server support."
                )

        except Exception as e:
            system_notes.append(f"❌ Failed to complete `{action_type}`: {e}")

    return files, memory, system_notes


# =============================================================================
# HTTP handler (Vercel Python runtime contract)
# =============================================================================

class handler(BaseHTTPRequestHandler):

    def _cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def do_OPTIONS(self):
        self.send_response(204)
        self._cors_headers()
        self.end_headers()

    def do_GET(self):
        self.send_response(200)
        self._cors_headers()
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps({"status": "Devgent backend is online"}).encode())

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            raw_body = self.rfile.read(length) if length else b"{}"
            payload = json.loads(raw_body or b"{}")

            message = (payload.get("message") or "").strip()
            history = payload.get("history") or []
            api_key = (payload.get("apiKey") or "").strip()
            model = payload.get("model") or "gemini-3.6-flash"
            memory = payload.get("memory") or {}
            attachment = payload.get("attachment")

            if not api_key:
                if message.strip().lower() == "help":
                    reply = GEMINI_API_KEY_HELP_TEXT
                else:
                    reply = "Please enter a free Gemini API key at the top of the page. Type 'help' if you want to know how to get one."
                return self._respond(200, {"reply": reply, "files": [], "memory": memory})

            if not message and not attachment:
                return self._respond(400, {"error": "Empty message."})

            full_prompt = message
            if attachment and attachment.get("data_base64"):
                try:
                    raw_bytes = base64.b64decode(attachment["data_base64"])
                    extracted = extract_file_content(attachment.get("filename", "attachment"), raw_bytes)
                    if extracted:
                        full_prompt += f"\n\n[Attached File Content ({attachment.get('filename', 'attachment')})]:\n{extracted[:20000]}"
                except Exception as e:
                    full_prompt += f"\n\n[Could not read attached file: {e}]"

            full_prompt += (
                "\n\n[System reminder: If you learned anything new about the user, the project, "
                "their preferences, or task status this turn, include a manage_memory action to save it - "
                "even if the user did not explicitly ask you to remember it.]"
            )

            memory_str = json.dumps(memory, indent=2) if memory else "No stored memories yet."
            sys_inst = SYSTEM_INSTRUCTION_TEMPLATE.format(current_date=date.today().isoformat(), memory_state=memory_str)

            response_text = call_gemini(api_key, model, sys_inst, history, full_prompt)
            response_text = re.sub(r'<think>[\s\S]*?</think>', '', response_text, flags=re.IGNORECASE).strip()

            # Pull the trailing ```json ... ``` action block out of the reply, same approach as the desktop app.
            json_match = re.search(r"```json\s*(\{[\s\S]*?\})\s*```", response_text)
            reply_text = response_text
            actions = []
            if json_match:
                reply_text = response_text[:json_match.start()].strip()
                try:
                    parsed = json.loads(json_match.group(1))
                    actions = parsed.get("actions", [])
                except Exception:
                    pass

            files, memory, system_notes = execute_actions(actions, memory) if actions else ([], memory, [])
            if system_notes:
                reply_text = (reply_text + "\n\n" + "\n".join(system_notes)).strip()

            self._respond(200, {"reply": reply_text, "files": files, "memory": memory})

        except Exception as e:
            traceback.print_exc()
            self._respond(500, {"error": str(e)})

    def _respond(self, status, obj):
        self.send_response(status)
        self._cors_headers()
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps(obj).encode("utf-8"))
