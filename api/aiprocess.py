"""
Devgent — Web Backend (Vercel Python Serverless Function)
===========================================================
Deployed at:  /api/AIProcess  (Vercel auto-routes api/AIProcess.py -> /api/AIProcess)
Called by:    chat.html

This is the browser-facing counterpart to the Devgent desktop app. Same job —
turn a plain-English request into a real Word doc / PowerPoint / spreadsheet /
PDF / code file — but re-architected for a stateless, public serverless
function instead of a trusted desktop process:

  * No `run_command` / `start_dev_server` actions. Those meant "execute an
    arbitrary shell command" on the desktop, where the *user* is the one
    running it on their own machine. Exposing that over a public HTTP
    endpoint would let anyone who finds this URL run arbitrary code on
    Vercel's infrastructure — a critical remote-code-execution hole. So the
    web version only ever *generates files* and hands them back as
    downloads; it never executes anything.
  * No server-side memory store. Vercel functions don't persist state
    between invocations, so "memory" is round-tripped with the client on
    every request instead of written to disk (see the `memory` field below).
  * The Gemini API key is supplied by the browser on every request and is
    never logged, stored, or written anywhere server-side.

Request body (POST, JSON):
{
  "message":  "user's chat message",
  "history":  [{"role": "user"|"model", "content": "..."}, ...]   # prior turns, most recent last, EXCLUDING this message
  "apiKey":   "user's Gemini API key",
  "model":    "gemini-3.6-flash" | "gemini-3.5-flash-lite",
  "memory":   { ...arbitrary persisted key/values... },
  "attachment": {"filename": "notes.docx", "data_base64": "..."},   # optional
  "connections": {                                                  # optional, one entry per connected tool
    "github":      {"access_token": "..."},
    "slack":       {"access_token": "..."},
    "notion":      {"access_token": "..."},
    "figma":       {"access_token": "..."},
    "microsoft365":{"access_token": "..."}
  }
}

Response body (JSON):
{
  "reply":  "assistant's chat text (markdown, JSON action block stripped)",
  "files":  [{"filename": "...", "mimetype": "...", "data_base64": "...", "group": "optional zip name if part of a batch"}],
  "memory": { ...updated memory dict, echo this back on the next request... },
  "connector_suggestions": [{"tool": "github", "reason": "so I can open this as a PR"}],
  "connectors": ["github", "slack", ...]   # echoes back which tools this request had a live connection for
}

=== THIRD-PARTY CONNECTORS (GitHub / Slack / Notion / Figma / Microsoft 365) ===
This backend doubles as an OAuth broker + API proxy for the tools people actually use
at work, so the model can act on their behalf once they've connected an account:

  * GET  /api/aiprocess.py?connect=<tool>
      Redirects the browser into that tool's own OAuth consent screen. <tool> is one
      of the keys in CONNECTORS below. The redirect_uri sent to the provider always
      points back at this same function.
  * GET  /api/aiprocess.py?oauth_callback=<tool>&code=...
      The OAuth provider redirects here after the person approves access. This
      exchanges the code for an access token *server-side* (using the app's client
      secret, which is only ever read from an env var and never sent to the
      browser), then returns a tiny HTML page that does
      `window.opener.postMessage(...)` with the resulting access token and closes
      itself. The browser tab that opened the popup is what actually stores the
      token (in localStorage) and sends it back on future POST requests inside
      "connections" — this backend never persists tokens anywhere, same trust model
      as the Gemini API key above.
  * Every POST request can include a "connections" object. Any tool with an
    "access_token" present is treated as connected for that turn, and the model is
    told so in its system instructions so it can actually call github_*, slack_*,
    notion_*, figma_*, or ms365_* actions instead of just asking the person to
    connect. See CONNECTORS / oauth_authorize_url / oauth_exchange_code below.
"""

from http.server import BaseHTTPRequestHandler
import json
import os
import re
import io
import zipfile
import base64
import traceback
from datetime import date

import requests

GEMINI_API_BASE = "https://generativelanguage.googleapis.com/v1beta/models"

GEMINI_API_KEY_HELP_TEXT = """Getting a Gemini API key is completely free and only takes a minute:

1. Open Google AI Studio in your browser at https://aistudio.google.com/apikey
2. Sign in with your Google account if you're asked to.
3. Click "Create API key", then choose to create a new project (or pick an existing one) when prompted.
4. Copy the API key that gets generated.
5. Paste it into the API key field at the top of the page.

No credit card or payment is required — the free tier is enough to chat with Devgent."""


# =============================================================================
# Third-party connectors — OAuth apps are all free to create; you (the site
# operator) create ONE app per tool, put its client id/secret in Vercel env
# vars, and then ANY visitor to your deployed site can click "Connect" and
# authorize their own account against your app. You never touch their
# credentials — the provider hands back a per-user access token.
# =============================================================================

CONNECTORS = {
    "github": {
        "label": "GitHub",
        "authorize_url": "https://github.com/login/oauth/authorize",
        "token_url": "https://github.com/login/oauth/access_token",
        "scope": "repo read:user",
        "client_id_env": "GITHUB_CLIENT_ID",
        "client_secret_env": "GITHUB_CLIENT_SECRET",
        "auth_style": "form",
        "setup_help": (
            "1. Go to https://github.com/settings/developers -> \"OAuth Apps\" -> \"New OAuth App\".\n"
            "2. Homepage URL: your deployed site's URL. Authorization callback URL: "
            "<your-site>/api/aiprocess.py?oauth_callback=github\n"
            "3. Free, no approval wait. Copy the Client ID, then click \"Generate a new client secret\" and copy that too.\n"
            "4. Set env vars GITHUB_CLIENT_ID and GITHUB_CLIENT_SECRET in Vercel to those two values.\n"
            "5. The `repo` scope already requested covers reading/committing/deleting files, branches, pull "
            "requests, and issues — no extra setup needed for those finer-grained actions."
        ),
    },
    "slack": {
        "label": "Slack",
        "authorize_url": "https://slack.com/oauth/v2/authorize",
        "token_url": "https://slack.com/api/oauth.v2.access",
        "scope": "chat:write channels:read channels:history groups:history reactions:write",
        "user_scope": "",
        "client_id_env": "SLACK_CLIENT_ID",
        "client_secret_env": "SLACK_CLIENT_SECRET",
        "auth_style": "form",
        "setup_help": (
            "1. Go to https://api.slack.com/apps -> \"Create New App\" -> \"From scratch\" (free, any Slack account).\n"
            "2. Under \"OAuth & Permissions\", add redirect URL: <your-site>/api/aiprocess.py?oauth_callback=slack\n"
            "3. Under \"Bot Token Scopes\" add: chat:write, channels:read, channels:history, groups:history, reactions:write "
            "(groups:history is needed to read private-channel messages; skip it if you only need public channels).\n"
            "4. Under \"Basic Information\" copy the Client ID and Client Secret.\n"
            "5. Set env vars SLACK_CLIENT_ID and SLACK_CLIENT_SECRET in Vercel."
        ),
    },
    "notion": {
        "label": "Notion",
        "authorize_url": "https://api.notion.com/v1/oauth/authorize",
        "token_url": "https://api.notion.com/v1/oauth/token",
        "scope": "",
        "client_id_env": "NOTION_CLIENT_ID",
        "client_secret_env": "NOTION_CLIENT_SECRET",
        "auth_style": "basic",
        "extra_authorize_params": {"owner": "user"},
        "setup_help": (
            "1. Go to https://www.notion.so/my-integrations -> \"New integration\".\n"
            "2. Set type to \"Public\" (this is what lets ANY Notion user connect, not just your own workspace).\n"
            "3. Redirect URI: <your-site>/api/aiprocess.py?oauth_callback=notion\n"
            "4. Under \"Capabilities\" enable Read content, Update content, and Insert content — without these, "
            "notion_read_page / notion_update_page / notion_create_page / notion_archive_page will fail even "
            "though the OAuth connection itself succeeds.\n"
            "5. Free. Copy the OAuth client ID and client secret from the integration's \"Distribution\" tab.\n"
            "6. Set env vars NOTION_CLIENT_ID and NOTION_CLIENT_SECRET in Vercel."
        ),
    },
    "figma": {
        "label": "Figma",
        "authorize_url": "https://www.figma.com/oauth",
        "token_url": "https://api.figma.com/v1/oauth/token",
        "scope": "files:read file_comments:write",
        "client_id_env": "FIGMA_CLIENT_ID",
        "client_secret_env": "FIGMA_CLIENT_SECRET",
        "auth_style": "form",
        "setup_help": (
            "1. Go to https://www.figma.com/developers/apps -> \"Create new app\" (free).\n"
            "2. Callback URL: <your-site>/api/aiprocess.py?oauth_callback=figma\n"
            "3. Copy the Client ID and Client secret.\n"
            "4. Set env vars FIGMA_CLIENT_ID and FIGMA_CLIENT_SECRET in Vercel.\n"
            "5. Figma's OAuth scope names occasionally change on their end (\"files:read\" vs older \"file_read\") — "
            "if figma_get_file/figma_export_images 401s, check the exact scope name Figma's dev portal shows for "
            "your app and update the `scope` value in CONNECTORS[\"figma\"] to match."
        ),
    },
    "microsoft365": {
        "label": "Microsoft 365",
        "authorize_url": "https://login.microsoftonline.com/common/oauth2/v2.0/authorize",
        "token_url": "https://login.microsoftonline.com/common/oauth2/v2.0/token",
        "scope": "offline_access Files.Read.All Sites.Read.All ChannelMessage.Read.All ChannelMessage.Send User.Read",
        "client_id_env": "MS365_CLIENT_ID",
        "client_secret_env": "MS365_CLIENT_SECRET",
        "auth_style": "form",
        "setup_help": (
            "1. Go to https://portal.azure.com -> \"Microsoft Entra ID\" -> \"App registrations\" -> \"New registration\" "
            "(a free Microsoft/Azure account is enough — no paid Azure subscription needed to register an app).\n"
            "2. Supported account types: \"Accounts in any organizational directory and personal Microsoft accounts\", "
            "so any visitor's work or personal Microsoft account can connect.\n"
            "3. Redirect URI (platform: Web): <your-site>/api/aiprocess.py?oauth_callback=microsoft365\n"
            "4. Under \"Certificates & secrets\" create a new client secret and copy its VALUE (not the ID).\n"
            "5. Under \"API permissions\" add the delegated Microsoft Graph permissions listed in \"scope\" above.\n"
            "   Files/Sites/User read scopes work immediately; ChannelMessage.Read.All and ChannelMessage.Send may "
            "prompt an admin-consent screen for work/school accounts — personal accounts and SharePoint/OneDrive "
            "access (recent files, search, reading a file's contents) work without it.\n"
            "6. Set env vars MS365_CLIENT_ID and MS365_CLIENT_SECRET in Vercel."
        ),
    },
}


def connected_tools(connections):
    """Which CONNECTORS keys have a usable access_token in this request's `connections` payload."""
    connections = connections or {}
    return [t for t in CONNECTORS if (connections.get(t) or {}).get("access_token")]


def oauth_authorize_url(tool, base_url, state=""):
    cfg = CONNECTORS[tool]
    client_id = os.environ.get(cfg["client_id_env"], "")
    params = {
        "client_id": client_id,
        "redirect_uri": f"{base_url}?oauth_callback={tool}",
        "response_type": "code",
        "state": state,
    }
    if cfg.get("scope"):
        params["scope"] = cfg["scope"]
    params.update(cfg.get("extra_authorize_params", {}))
    if not client_id:
        return None
    from urllib.parse import urlencode
    return f"{cfg['authorize_url']}?{urlencode(params)}"


def oauth_exchange_code(tool, code, base_url):
    """Server-side code -> access_token exchange. The client secret never leaves this function."""
    cfg = CONNECTORS[tool]
    client_id = os.environ.get(cfg["client_id_env"], "")
    client_secret = os.environ.get(cfg["client_secret_env"], "")
    if not client_id or not client_secret:
        raise RuntimeError(
            f"{cfg['label']} isn't configured on this deployment yet — "
            f"{cfg['client_id_env']} / {cfg['client_secret_env']} env vars are missing."
        )

    data = {
        "client_id": client_id,
        "client_secret": client_secret,
        "code": code,
        "redirect_uri": f"{base_url}?oauth_callback={tool}",
        "grant_type": "authorization_code",
    }
    headers = {"Accept": "application/json"}
    if cfg.get("auth_style") == "basic":
        resp = requests.post(cfg["token_url"], data=data, headers=headers,
                              auth=(client_id, client_secret), timeout=30)
    else:
        resp = requests.post(cfg["token_url"], data=data, headers=headers, timeout=30)

    try:
        payload = resp.json()
    except Exception:
        raise RuntimeError(f"{cfg['label']} token exchange returned a non-JSON response ({resp.status_code}).")

    token = payload.get("access_token")
    if not token:
        raise RuntimeError(f"{cfg['label']} token exchange failed: {payload.get('error_description') or payload.get('error') or payload}")
    return payload


SYSTEM_INSTRUCTION_TEMPLATE = """
You are Devgent, an AI assistant running as a web app in the person's browser.
Current date: {current_date}

=== CORE OPERATING PRINCIPLES ===
1. Respond clearly to the user in plain chat text, then, ONLY if you are actually taking an action this turn, attach a single JSON block at the very end of your response inside triple backticks tagged with json.
2. The JSON block can contain an "actions" list so you can queue multiple actions in a single turn.
3. File Naming: Always specify exact file extensions (.py, .js, .html, .css, .json, .txt, .pdf, .docx, .xlsx, .pptx).
4. Environment Awareness: You are running as a stateless web backend. You have NO terminal access and CANNOT run commands, install packages, or start servers. You can only generate files, which are handed to the user as downloads in their browser. Never claim to run or execute code — only write it.
5. No Permission Theater Needed: Unlike a desktop app, generating a file here is harmless (nothing on the user's machine is touched until they click Download), so you do not need to ask "can I create this file?" before emitting a create_* action — just briefly describe what you're making, then issue the action.
6. IMPORTANT — Never claim a file has been "saved" or "created on your computer": this is a stateless web backend with no access to the user's file system. When a create_* action finishes, a download card with a Download button appears automatically in the chat — you do not need to mention how to download it, but you must phrase your own reply as "I've put together..." / "Here's your..." / "Ready below —", never "I've saved this to your Documents folder" or any wording implying direct disk access.
7. Markdown In Chat: Your plain-text reply (everything outside the JSON block) is rendered as real Markdown — headings, **bold**, *italics*, bullet/numbered lists, tables, and fenced code blocks with a language tag (e.g. ```python ... ```) will all display properly, including syntax-highlighted code blocks. This is completely separate from the JSON action block used to trigger file creation.

=== CLARIFY BEFORE ACTING (VERY IMPORTANT) ===
If the user's request is vague, incomplete, or does not tell you concretely WHAT to build, write, or do — for example "I want to make a python script", "build me something cool", "make a presentation", "help me with a document" — then you must NOT take any action and must NOT emit a JSON action block at all this turn. Instead:
- Ask the user, in plain chat text only, exactly what they want: purpose, subject matter, desired content/sections, style, filename, etc. — whatever is missing.
- If they already gave you some partial detail, save it to memory with "manage_memory" so it isn't lost (this is the one exception where you may still emit a JSON block even without creating a file), e.g.:
```json
{{
  "actions": [
    {{"action": "manage_memory", "action_data": {{"updates": {{"pending_request": "User wants a Python script, no further details yet."}}}}}}
  ]
}}
```
- Wait for the user's next message before creating anything.
Only move on to create_txt / create_docx / create_pdf / create_pptx / create_xlsx / create_batch_files once the user has given you enough concrete specifics to act on.

=== CONNECTED TOOLS (GitHub, Slack, Notion, Figma, Microsoft 365) ===
Devgent can connect to the tools office workers and developers actually use, and act
on their behalf once they've authorized it. Status for THIS person, right now:
{connectors_state}

Rules for using this:
1. If a tool is listed as CONNECTED, you may actually use its action(s) below when it
   would help — e.g. push generated code to a GitHub repo, post a message to Slack,
   create a Notion page, pull Figma design context, or read SharePoint/Teams context —
   instead of just generating a local file. Ask for any missing specifics you need
   (repo name, channel, page, file key) the same way you'd clarify anything else.
2. If a tool is NOT connected and the user's message names it, references it, or
   clearly implies it (e.g. "push this to my repo", "post this in our channel",
   "put this in Notion", "check the Figma file", "what does the SharePoint doc say")
   — include a `suggest_connector` action for that tool THIS TURN, in addition to
   your normal reply/actions. Do this every time it's relevant, not just once.
3. Routinely (roughly every 3-5 turns, when it hasn't come up on its own) suggest
   connecting ONE specific unconnected tool that plausibly fits what the person is
   doing — name the exact tool, never a vague "connect a tool related to your work".
   Skip this if every tool is already connected, or if the conversation is too short
   to have a good guess yet.
4. Never claim you performed a connected-tool action you didn't actually emit an
   action for, and never claim a tool is connected when the status above says it isn't.

=== DYNAMIC MEMORY STORE (BE AGGRESSIVE ABOUT USING THIS) ===
Current Stored Memories:
{memory_state}

You can save, update, or discard memory items dynamically to retain user context across conversations. Treat this as a first-class habit — on almost every turn, ask yourself "did I just learn something worth remembering?" That includes, but is not limited to:
- The user's name, role, or how they prefer to be addressed.
- Project name, purpose, and key details already established.
- Stylistic or technical preferences (tone, color themes, frameworks).
- Decisions the user has made or corrections they've given you (these should overwrite/replace old memory values, not just add new ones).
- The current status of a multi-step task, so you can resume correctly next session.
If the answer is yes, include a "manage_memory" action alongside any other action in the same "actions" array. Err on the side of saving too much rather than too little.

=== MULTI-ACTION SCHEMA ===
```json
{{
  "actions": [
    {{"action": "create_txt", "action_data": {{"filename": "app.py", "content": "print('hello')"}}}}
  ]
}}
```

=== AVAILABLE ACTIONS ===

1. Manage Dynamic Memory (manage_memory)
{{"action": "manage_memory", "action_data": {{"updates": {{"user_theme": "dark"}}, "discards": ["old_project_key"]}}}}

2. Advanced Styled PowerPoint Presentation (create_pptx)
```json
{{
  "action": "create_pptx",
  "action_data": {{
    "filename": "pitch.pptx",
    "title": "Deck Title",
    "subtitle": "Optional subtitle line",
    "title_color": "#3B82F6",
    "subtitle_color": "#94A3B8",
    "background_color": "#0F172A",
    "theme_accent": "#3B82F6",
    "slides": [
      {{
        "layout": "bullets",
        "title": "Slide Heading",
        "background_color": "#F8FAFC",
        "header_color": "#1E293B",
        "header_text_color": "#FFFFFF",
        "accent_color": "#3B82F6",
        "content": ["First point", "Second point", "Third point"],
        "chart": {{
          "type": "column",
          "categories": ["Q1", "Q2", "Q3"],
          "series": [{{"name": "Revenue", "values": [10, 20, 30]}}]
        }}
      }},
      {{
        "layout": "two_column",
        "title": "Comparison",
        "left_heading": "Pros",
        "content_left": ["Fast", "Cheap"],
        "right_heading": "Cons",
        "content_right": ["Riskier"]
      }},
      {{
        "layout": "quote",
        "title": "Customer Voice",
        "quote_text": "This product changed how we work.",
        "quote_author": "Jane Doe, CTO"
      }},
      {{
        "layout": "section",
        "title": "Part Two: Results",
        "subtitle": "Optional divider subtitle",
        "background_color": "#3B82F6",
        "header_text_color": "#FFFFFF"
      }}
    ]
  }}
}}
```
Rules:
- "layout" is one of: "bullets" (default; supports an optional "chart" side-by-side), "two_column", "quote", or "section" (a full-bleed divider slide).
- All colors are hex strings like "#0F172A". Any slide can override background_color, header_color, header_text_color, and accent_color; otherwise it inherits "theme_accent".
- Chart "type" is one of "bar", "column", "pie", "line".
- Prefer 3-6 short bullets per slide rather than long paragraphs.

3. Excel Spreadsheet (create_xlsx)
{{"action": "create_xlsx", "action_data": {{"filename": "data.xlsx", "sheet_name": "Metrics", "headers": ["Name", "Score"], "rows": [["Alice", 95]]}}}}

4. Word Document (create_docx)
{{"action": "create_docx", "action_data": {{"filename": "doc.docx", "title": "Proposal", "heading": "Summary", "paragraphs": ["Details..."]}}}}

5. Styled PDF Document (create_pdf)
```json
{{
  "action": "create_pdf",
  "action_data": {{
    "filename": "report.pdf",
    "title": "Report Title",
    "subtitle": "Optional subtitle",
    "accent_color": "#3B82F6",
    "sections": [
      {{
        "heading": "Overview",
        "paragraphs": ["Body paragraph text goes here."],
        "bullets": ["Point one", "Point two"],
        "quote": "An optional pull-quote highlighting something important.",
        "table": {{"headers": ["Name", "Score"], "rows": [["Alice", "95"], ["Bob", "88"]]}}
      }}
    ]
  }}
}}
```
- For very simple documents you may skip "sections" and just pass a flat "paragraphs" list.

6. Single Text/Code File (create_txt)
{{"action": "create_txt", "action_data": {{"filename": "app.py", "content": "print('hello')"}}}}

7. Batch Files, packaged as one .zip (create_batch_files)
{{"action": "create_batch_files", "action_data": {{"zip_filename": "site.zip", "files": [{{"filename": "index.html", "content": "<h1>Hello</h1>"}}]}}}}
Note: when files reference each other (an .html that loads a sibling .js/.css from this same batch), use plain relative filenames
("script.js", not "./script.js" or an absolute/offsite URL) so the in-chat live preview can wire them together.

8. Suggest Connecting A Tool (suggest_connector) — use for any NOT-yet-connected tool per the rules above
{{"action": "suggest_connector", "action_data": {{"tool": "github", "reason": "so I can commit this straight to your repo"}}}}
"tool" must be one of: github, slack, notion, figma, microsoft365.

9. GitHub (only when connected)
{{"action": "github_list_repos", "action_data": {{}}}}
{{"action": "github_list_files", "action_data": {{"owner": "octocat", "repo": "my-site", "path": "src", "ref": "optional-branch-or-sha"}}}}
{{"action": "github_read_file", "action_data": {{"owner": "octocat", "repo": "my-site", "path": "src/app.py", "ref": "optional-branch-or-sha"}}}}
{{"action": "github_commit_file", "action_data": {{"owner": "octocat", "repo": "my-site", "path": "src/app.py", "content": "print('hi')", "message": "Devgent: update app.py", "branch": "main"}}}}
{{"action": "github_delete_file", "action_data": {{"owner": "octocat", "repo": "my-site", "path": "old.py", "message": "Devgent: remove old.py", "branch": "main"}}}}
{{"action": "github_create_branch", "action_data": {{"owner": "octocat", "repo": "my-site", "branch": "devgent-feature", "base": "main"}}}}
{{"action": "github_list_pull_requests", "action_data": {{"owner": "octocat", "repo": "my-site", "state": "open"}}}}
{{"action": "github_create_pull_request", "action_data": {{"owner": "octocat", "repo": "my-site", "title": "Add feature", "head": "devgent-feature", "base": "main", "body": "What changed and why."}}}}
{{"action": "github_merge_pull_request", "action_data": {{"owner": "octocat", "repo": "my-site", "pr_number": 12, "merge_method": "merge"}}}}
{{"action": "github_close_pull_request", "action_data": {{"owner": "octocat", "repo": "my-site", "pr_number": 12}}}}
{{"action": "github_list_issues", "action_data": {{"owner": "octocat", "repo": "my-site", "state": "open"}}}}
{{"action": "github_create_issue", "action_data": {{"owner": "octocat", "repo": "my-site", "title": "Bug: X breaks", "body": "Repro steps..."}}}}
{{"action": "github_comment_on_issue", "action_data": {{"owner": "octocat", "repo": "my-site", "issue_number": 4, "body": "Comment text"}}}}
{{"action": "github_close_issue", "action_data": {{"owner": "octocat", "repo": "my-site", "issue_number": 4}}}}

10. Slack (only when connected)
{{"action": "slack_list_channels", "action_data": {{}}}}
{{"action": "slack_read_channel_messages", "action_data": {{"channel": "C0123456", "limit": 20}}}}
{{"action": "slack_send_message", "action_data": {{"channel": "#general", "text": "Message text"}}}}
{{"action": "slack_reply_thread", "action_data": {{"channel": "#general", "thread_ts": "1699999999.000200", "text": "Reply text"}}}}
{{"action": "slack_add_reaction", "action_data": {{"channel": "#general", "timestamp": "1699999999.000200", "emoji": "tada"}}}}
Note: slack_list_channels returns each channel's id — use that id (not the #name) for read/reply/react actions.

11. Notion (only when connected)
{{"action": "notion_search", "action_data": {{"query": "roadmap"}}}}
{{"action": "notion_read_page", "action_data": {{"page_id": "id-from-a-prior-search"}}}}
{{"action": "notion_create_page", "action_data": {{"parent_page_id": "optional-id-from-a-prior-search", "title": "Meeting Notes", "paragraphs": ["First paragraph.", "Second paragraph."]}}}}
{{"action": "notion_update_page", "action_data": {{"page_id": "id-from-a-prior-search", "paragraphs": ["New paragraph appended to the page."]}}}}
{{"action": "notion_archive_page", "action_data": {{"page_id": "id-from-a-prior-search"}}}}

12. Figma (only when connected)
{{"action": "figma_get_file", "action_data": {{"file_key": "abc123"}}}}
{{"action": "figma_list_comments", "action_data": {{"file_key": "abc123"}}}}
{{"action": "figma_add_comment", "action_data": {{"file_key": "abc123", "message": "Comment text", "node_id": "optional-1:2"}}}}
{{"action": "figma_export_images", "action_data": {{"file_key": "abc123", "node_ids": ["1:2", "1:3"], "format": "png"}}}}
The file_key is the id segment in a Figma file URL: figma.com/file/<file_key>/...

13. Microsoft 365 (only when connected)
{{"action": "ms365_recent_files", "action_data": {{}}}}
{{"action": "ms365_search_files", "action_data": {{"query": "Q3 budget"}}}}
{{"action": "ms365_read_file", "action_data": {{"item_id": "id-from-recent-or-search"}}}}
{{"action": "ms365_list_teams_channels", "action_data": {{"team_id": "optional — omit to list joined teams instead"}}}}
{{"action": "ms365_teams_messages", "action_data": {{"team_id": "...", "channel_id": "..."}}}}
{{"action": "ms365_send_teams_message", "action_data": {{"team_id": "...", "channel_id": "...", "text": "Message text"}}}}
(Use ms365_recent_files/search_files or ms365_list_teams_channels first if you don't already have an item_id/team_id/channel_id in memory or history.)

=== CAUTION WITH HARD-TO-UNDO CONNECTOR ACTIONS ===
merge_pull_request, close_pull_request, delete_file, close_issue, and notion_archive_page change or
remove something real for the person and are not easily reversible from here. Only take one of these
when the user's message unambiguously asks for exactly that action on that specific target — never
chain one onto a broader request on your own initiative (e.g. don't close an issue just because you
opened a PR that references it, unless asked to).

Remember to USE your memory as much as you can. You will remember absolutely nothing that you do not store in your memory.
Furthermore, if the user asks you to look at a specific file in their Github repository, open the repository, then locate the file they asked you to look at, then open the file. Do not simply just check the repositories you can open, or just open the repository but not look at the specific file. Likewise, do not do similar stuff for Notion, Figma, Slack, or Microsoft 365
"""


# =============================================================================
# Small color helpers (shared by docx/pptx/pdf/xlsx builders)
# =============================================================================

def hex_to_rgb_tuple(hex_str):
    hex_str = (hex_str or "").lstrip('#')
    if len(hex_str) != 6:
        hex_str = "3B82F6"
    return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)


def shade_color(hex_str, amount):
    r, g, b = hex_to_rgb_tuple(hex_str)
    if amount >= 0:
        r = int(r + (255 - r) * amount)
        g = int(g + (255 - g) * amount)
        b = int(b + (255 - b) * amount)
    else:
        r, g, b = int(r * (1 + amount)), int(g * (1 + amount)), int(b * (1 + amount))
    r, g, b = max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b))
    return f"#{r:02X}{g:02X}{b:02X}"


def safe_filename(raw, default_ext):
    name = os.path.basename((raw or "").strip() or f"file{default_ext}")
    name = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "_", name)
    if not name.lower().endswith(default_ext):
        name += default_ext
    return name


# =============================================================================
# File builders — each returns raw bytes
# =============================================================================

def build_docx_bytes(data):
    import docx
    from docx.shared import Pt, RGBColor, Inches, Emu
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    doc = docx.Document()
    title_text = data.get("title", "Document Title")
    p_title = doc.add_paragraph()
    run = p_title.add_run(title_text)
    run.font.name = 'Calibri'
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

    accent_table = doc.add_table(rows=1, cols=1)
    accent_cell = accent_table.rows[0].cells[0]
    accent_cell.width = Inches(6.0)
    shading = OxmlElement('w:shd')
    shading.set(qn('w:fill'), "3B82F6")
    accent_cell._tc.get_or_add_tcPr().append(shading)
    accent_table.rows[0].height = Emu(9525 * 3)

    if data.get("heading"):
        p_head = doc.add_paragraph()
        r_head = p_head.add_run(data["heading"])
        r_head.font.name = 'Calibri'
        r_head.font.size = Pt(16)
        r_head.font.bold = True
        r_head.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)

    for ptext in data.get("paragraphs", []):
        p = doc.add_paragraph()
        r = p.add_run(ptext)
        r.font.name = 'Calibri'
        r.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pdf_bytes(data):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    )

    accent_hex = data.get("accent_color", "#3B82F6")
    accent = colors.HexColor(accent_hex)
    accent_dark = colors.HexColor(shade_color(accent_hex, -0.35))
    text_dark = colors.HexColor("#1F2937")
    muted = colors.HexColor("#64748B")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter,
        topMargin=0.9 * inch, bottomMargin=0.8 * inch,
        leftMargin=0.75 * inch, rightMargin=0.75 * inch
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('T', parent=styles['Heading1'], fontSize=26, leading=30, textColor=colors.white, spaceAfter=4, fontName="Helvetica-Bold")
    subtitle_style = ParagraphStyle('S', parent=styles['Normal'], fontSize=12, leading=16, textColor=colors.HexColor("#E2E8F0"))
    heading_style = ParagraphStyle('H', parent=styles['Heading2'], fontSize=15, leading=19, textColor=accent_dark, spaceBefore=18, spaceAfter=8, fontName="Helvetica-Bold")
    body_style = ParagraphStyle('B', parent=styles['Normal'], fontSize=10.5, leading=15, textColor=text_dark, spaceAfter=8)
    bullet_style = ParagraphStyle('BL', parent=body_style, leftIndent=14, spaceAfter=6)
    quote_style = ParagraphStyle('Q', parent=body_style, fontName="Helvetica-Oblique", leftIndent=16, textColor=muted)

    story = []
    title_text = data.get("title", "Generated Report")
    subtitle_text = data.get("subtitle", "")
    title_cell = [Paragraph(title_text, title_style)]
    if subtitle_text:
        title_cell.append(Paragraph(subtitle_text, subtitle_style))
    title_table = Table([[title_cell]], colWidths=[6.9 * inch])
    title_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), accent),
        ('LEFTPADDING', (0, 0), (-1, -1), 18), ('RIGHTPADDING', (0, 0), (-1, -1), 18),
        ('TOPPADDING', (0, 0), (-1, -1), 18), ('BOTTOMPADDING', (0, 0), (-1, -1), 18),
    ]))
    story.append(title_table)
    story.append(Spacer(1, 18))

    sections = data.get("sections") or [{"paragraphs": data.get("paragraphs", [])}]
    for section in sections:
        if section.get("heading"):
            h_table = Table([[Paragraph(section["heading"], heading_style)]], colWidths=[6.9 * inch])
            h_table.setStyle(TableStyle([
                ('LINEBELOW', (0, 0), (-1, -1), 1.4, accent),
                ('LEFTPADDING', (0, 0), (-1, -1), 0), ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            story.append(h_table)

        for paragraph in section.get("paragraphs", []):
            story.append(Paragraph(paragraph, body_style))

        bullets = section.get("bullets", [])
        if bullets:
            items = [ListItem(Paragraph(b, bullet_style), bulletColor=accent, value='●') for b in bullets]
            story.append(ListFlowable(items, bulletType='bullet', start='●', leftIndent=16))
            story.append(Spacer(1, 6))

        if section.get("quote"):
            q_table = Table([[Paragraph(f"\u201c{section['quote']}\u201d", quote_style)]], colWidths=[6.9 * inch])
            q_table.setStyle(TableStyle([
                ('LINEBEFORE', (0, 0), (0, -1), 3, accent),
                ('LEFTPADDING', (0, 0), (-1, -1), 14),
                ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ]))
            story.append(q_table)
            story.append(Spacer(1, 8))

        table_data = section.get("table")
        if table_data and table_data.get("headers") and table_data.get("rows"):
            tbl = Table([table_data["headers"]] + table_data["rows"], hAlign='LEFT')
            tbl.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), accent),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9.5),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F5F9")]),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 10))
        story.append(Spacer(1, 4))

    def _footer(canvas, doc_):
        canvas.saveState()
        canvas.setStrokeColor(accent)
        canvas.setLineWidth(1)
        canvas.line(0.75 * inch, 0.6 * inch, letter[0] - 0.75 * inch, 0.6 * inch)
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(muted)
        canvas.drawString(0.75 * inch, 0.45 * inch, title_text)
        canvas.drawRightString(letter[0] - 0.75 * inch, 0.45 * inch, f"Page {doc_.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=_footer, onLaterPages=_footer)
    return buf.getvalue()


def build_xlsx_bytes(data):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = data.get("sheet_name", "Sheet1")

    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    alt_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
    thin_border = Border(*(Side(style='thin', color='CBD5E1') for _ in range(4)))

    headers = data.get("headers", [])
    if headers:
        ws.append(headers)
        for col_num in range(1, len(headers) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

    rows = data.get("rows", [])
    for row_idx, row in enumerate(rows, start=2):
        ws.append(row)
        for col_num in range(1, len(row) + 1):
            cell = ws.cell(row=row_idx, column=col_num)
            cell.font = Font(name="Calibri", size=10)
            cell.border = thin_border
            if row_idx % 2 == 0:
                cell.fill = alt_fill

    for col in ws.columns:
        max_len = max((len(str(c.value or '')) for c in col), default=0)
        ws.column_dimensions[get_column_letter(col[0].column)].width = max(max_len + 4, 12)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _add_pptx_footer(slide, index, total, accent_color, slide_w, slide_h, on_light_bg=True):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), slide_h - Inches(0.45), slide_w - Inches(1.0), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(accent_color)
    line.line.fill.background()
    line.shadow.inherit = False

    tb = slide.shapes.add_textbox(slide_w - Inches(1.6), slide_h - Inches(0.42), Inches(1.1), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{index:02d} / {total:02d}"
    p.font.size = Pt(10)
    p.font.color.rgb = _rgb(accent_color if on_light_bg else "#94A3B8")
    p.alignment = PP_ALIGN.RIGHT


def _rgb(hex_str):
    from pptx.dml.color import RGBColor
    r, g, b = hex_to_rgb_tuple(hex_str)
    return RGBColor(r, g, b)


def build_pptx_bytes(data):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
    theme_accent = data.get("theme_accent", data.get("title_color", "#3B82F6"))

    # ---- Title slide ----
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bg_hex = data.get("background_color", "#0F172A")
    bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = _rgb(title_bg_hex); bg.line.fill.background(); bg.shadow.inherit = False

    accent_circle = title_slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(4.0), Inches(-2.0), Inches(6.0), Inches(6.0))
    accent_circle.fill.solid(); accent_circle.fill.fore_color.rgb = _rgb(shade_color(theme_accent, -0.35))
    accent_circle.line.fill.background(); accent_circle.shadow.inherit = False

    accent_bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), Inches(0.18), Inches(1.9))
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = _rgb(theme_accent)
    accent_bar.line.fill.background(); accent_bar.shadow.inherit = False

    txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(9.5), Inches(3.0))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Presentation Title")
    p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = _rgb(data.get("title_color", theme_accent))

    sub = tf.add_paragraph()
    sub.text = data.get("subtitle", "Generated by Devgent")
    sub.font.size = Pt(20)
    sub.font.color.rgb = _rgb(data.get("subtitle_color", "#94A3B8"))
    sub.space_before = Pt(14)

    slides_list = data.get("slides", [])
    total_slides = len(slides_list)

    for i, sdata in enumerate(slides_list, start=1):
        layout = str(sdata.get("layout", "bullets")).lower()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        accent_color = sdata.get("accent_color", theme_accent)
        slide_bg_hex = sdata.get("background_color", "#3B82F6" if layout == "section" else "#F8FAFC")

        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        s_bg.fill.solid(); s_bg.fill.fore_color.rgb = _rgb(slide_bg_hex)
        s_bg.line.fill.background(); s_bg.shadow.inherit = False

        is_dark_bg = slide_bg_hex.upper() not in ["#FFFFFF", "#F8FAFC", "#F1F5F9", "#FFF"]
        card_fill = "#1E293B" if is_dark_bg else "#FFFFFF"
        default_text = sdata.get("text_color", "#F8FAFC" if is_dark_bg else "#334155")

        if layout == "section":
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.6))
            tf2 = tb.text_frame; tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sdata.get("title", "Section")
            p2.font.size = Pt(42); p2.font.bold = True
            p2.font.color.rgb = _rgb(sdata.get("header_text_color", "#FFFFFF"))
            p2.alignment = PP_ALIGN.CENTER
            if sdata.get("subtitle"):
                sp = tf2.add_paragraph()
                sp.text = sdata["subtitle"]
                sp.font.size = Pt(18)
                sp.font.color.rgb = _rgb(shade_color(sdata.get("header_text_color", "#FFFFFF"), -0.2))
                sp.alignment = PP_ALIGN.CENTER
            _add_pptx_footer(slide, i, total_slides, "#FFFFFF", SLIDE_W, SLIDE_H, on_light_bg=False)
            continue

        header_bg_hex = sdata.get("header_color", accent_color)
        header_txt_hex = sdata.get("header_text_color", "#FFFFFF")
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), SLIDE_W - Inches(1.0), Inches(0.9))
        banner.fill.solid(); banner.fill.fore_color.rgb = _rgb(header_bg_hex)
        banner.line.fill.background(); banner.shadow.inherit = False
        bp = banner.text_frame.paragraphs[0]
        bp.text = sdata.get("title", "Slide Title")
        bp.font.size = Pt(24); bp.font.bold = True
        bp.font.color.rgb = _rgb(header_txt_hex)

        if layout == "quote":
            quote_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.8), SLIDE_W - Inches(2.8), Inches(4.6))
            quote_card.fill.solid(); quote_card.fill.fore_color.rgb = _rgb(card_fill)
            quote_card.line.fill.background(); quote_card.shadow.inherit = False

            q_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(1.8), Inches(0.12), Inches(4.6))
            q_bar.fill.solid(); q_bar.fill.fore_color.rgb = _rgb(accent_color)
            q_bar.line.fill.background(); q_bar.shadow.inherit = False

            qtf = quote_card.text_frame; qtf.word_wrap = True
            qtf.margin_left = Inches(0.5); qtf.margin_right = Inches(0.4)
            qp = qtf.paragraphs[0]
            qp.text = f"\u201c{sdata.get('quote_text', '')}\u201d"
            qp.font.size = Pt(24); qp.font.italic = True
            qp.font.color.rgb = _rgb(default_text)
            if sdata.get("quote_author"):
                ap = qtf.add_paragraph()
                ap.text = f"\u2014 {sdata['quote_author']}"
                ap.font.size = Pt(16); ap.font.bold = True
                ap.font.color.rgb = _rgb(accent_color)
                ap.space_before = Pt(16)

        elif layout == "two_column":
            col_w = int((SLIDE_W - Inches(1.4)) / 2)
            columns = [
                (Inches(0.5), sdata.get("left_heading", ""), sdata.get("content_left", [])),
                (Inches(0.5) + col_w + Inches(0.4), sdata.get("right_heading", ""), sdata.get("content_right", [])),
            ]
            for left, heading, items in columns:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), col_w, Inches(5.4))
                card.fill.solid(); card.fill.fore_color.rgb = _rgb(card_fill)
                card.line.color.rgb = _rgb(shade_color(card_fill, -0.05))
                card.shadow.inherit = False
                ctf = card.text_frame; ctf.word_wrap = True
                ctf.margin_left = Inches(0.3); ctf.margin_top = Inches(0.25)
                first = False
                if heading:
                    hp = ctf.paragraphs[0]
                    hp.text = heading; hp.font.bold = True; hp.font.size = Pt(16)
                    hp.font.color.rgb = _rgb(accent_color); hp.space_after = Pt(10)
                    first = True
                for bullet in items:
                    bp2 = ctf.add_paragraph() if first else ctf.paragraphs[0]
                    first = True
                    r1 = bp2.add_run(); r1.text = "\u25cf "; r1.font.color.rgb = _rgb(accent_color); r1.font.size = Pt(14)
                    r2 = bp2.add_run(); r2.text = bullet; r2.font.color.rgb = _rgb(default_text); r2.font.size = Pt(14)
                    bp2.space_after = Pt(8)

        else:  # bullets (default), optionally with a chart
            content_items = sdata.get("content", [])
            has_chart = isinstance(sdata.get("chart"), dict)
            text_width = Inches(5.6) if has_chart else SLIDE_W - Inches(1.0)
            chart_left = Inches(6.5)
            chart_width = SLIDE_W - chart_left - Inches(0.5)

            if content_items:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), text_width, Inches(5.4))
                card.fill.solid(); card.fill.fore_color.rgb = _rgb(card_fill)
                card.line.color.rgb = _rgb(shade_color(card_fill, -0.05))
                card.shadow.inherit = False
                ctf = card.text_frame; ctf.word_wrap = True
                ctf.margin_left = Inches(0.35); ctf.margin_top = Inches(0.3)
                for idx, bullet in enumerate(content_items):
                    cp = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
                    r1 = cp.add_run(); r1.text = "\u25cf "; r1.font.size = Pt(15); r1.font.color.rgb = _rgb(accent_color)
                    r2 = cp.add_run(); r2.text = bullet; r2.font.size = Pt(15); r2.font.color.rgb = _rgb(default_text)
                    cp.space_after = Pt(10)

            if has_chart:
                chart_info = sdata["chart"]
                chart_data = CategoryChartData()
                chart_data.categories = chart_info.get("categories", [])
                for s in chart_info.get("series", []):
                    chart_data.add_series(s.get("name", "Series"), tuple(s.get("values", [])))
                chart_type_map = {
                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                    "pie": XL_CHART_TYPE.PIE, "line": XL_CHART_TYPE.LINE
                }
                xl_type = chart_type_map.get(str(chart_info.get("type", "column")).lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
                chart_shape = slide.shapes.add_chart(xl_type, chart_left, Inches(1.5), chart_width, Inches(5.4), chart_data)
                chart = chart_shape.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

        _add_pptx_footer(slide, i, total_slides, accent_color, SLIDE_W, SLIDE_H, on_light_bg=not is_dark_bg)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_batch_zip_bytes(files):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            fname = safe_filename(f.get("filename", "file.txt"), "")
            zf.writestr(fname, f.get("content", ""))
    return buf.getvalue()


# =============================================================================
# Attachment text extraction (mirrors the desktop app's extract_file_content)
# =============================================================================

def extract_file_content(filename, raw_bytes):
    ext = os.path.splitext(filename or "")[1].lower()
    try:
        if ext == '.docx':
            import docx
            doc = docx.Document(io.BytesIO(raw_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    if any(c is not None for c in row):
                        content.append("\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(content)
        elif ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        else:
            return raw_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[Error extracting file content from {filename}: {e}]"


# =============================================================================
# Gemini call (plain REST — no SDK, keeps the deployed function small)
# =============================================================================

def call_gemini(api_key, model, system_instruction, history, user_message):
    contents = []
    for h in history:
        role = "user" if h.get("role") == "user" else "model"
        contents.append({"role": role, "parts": [{"text": h.get("content", "")}]})
    contents.append({"role": "user", "parts": [{"text": user_message}]})

    body = {
        "contents": contents,
        "systemInstruction": {"parts": [{"text": system_instruction}]},
        "generationConfig": {"temperature": 0.3},
    }
    url = f"{GEMINI_API_BASE}/{model}:generateContent?key={api_key}"
    resp = requests.post(url, json=body, timeout=120)
    if resp.status_code != 200:
        detail = ""
        try:
            detail = resp.json().get("error", {}).get("message", "")
        except Exception:
            detail = resp.text[:300]
        raise RuntimeError(f"Gemini API error ({resp.status_code}): {detail}")

    data = resp.json()
    candidates = data.get("candidates") or []
    if not candidates:
        raise RuntimeError("Gemini returned no candidates (the response may have been blocked by safety filters).")
    parts = candidates[0].get("content", {}).get("parts", [])
    return "".join(p.get("text", "") for p in parts)


# =============================================================================
# Connector API calls — each takes the per-request access token the browser
# supplied in `connections`; nothing here is ever persisted server-side.
# =============================================================================

def _connector_token(connections, tool):
    token = ((connections or {}).get(tool) or {}).get("access_token")
    if not token:
        raise RuntimeError(f"{CONNECTORS[tool]['label']} isn't connected yet.")
    return token


def github_api(token, method, path, **kwargs):
    resp = requests.request(method, f"https://api.github.com{path}",
                             headers={"Authorization": f"Bearer {token}",
                                      "Accept": "application/vnd.github+json"},
                             timeout=30, **kwargs)
    if resp.status_code >= 400:
        raise RuntimeError(f"GitHub API error ({resp.status_code}): {resp.text[:300]}")
    return resp.json() if resp.text else {}


def do_github_list_repos(token, action_data, files_out):
    repos = github_api(token, "GET", "/user/repos?per_page=20&sort=updated")
    names = ", ".join(r["full_name"] for r in repos[:20]) or "(no repos found)"
    return f"Repos I can see: {names}"


def do_github_list_files(token, action_data, files_out):
    owner, repo = action_data["owner"], action_data["repo"]
    path = action_data.get("path", "")
    ref_q = f"?ref={action_data['ref']}" if action_data.get("ref") else ""
    items = github_api(token, "GET", f"/repos/{owner}/{repo}/contents/{path}{ref_q}")
    if isinstance(items, dict):  # a single file path was given, not a directory
        return f"`{path}` is a file, not a directory — use github_read_file to read it."
    lines = [f"- {it['path']}{'/' if it['type'] == 'dir' else ''}" for it in items]
    return f"Contents of {owner}/{repo}/{path or '(root)'}:\n" + ("\n".join(lines) if lines else "(empty)")


def do_github_read_file(token, action_data, files_out):
    owner, repo, path = action_data["owner"], action_data["repo"], action_data["path"]
    ref_q = f"?ref={action_data['ref']}" if action_data.get("ref") else ""
    result = github_api(token, "GET", f"/repos/{owner}/{repo}/contents/{path}{ref_q}")
    raw = base64.b64decode(result.get("content", "").replace("\n", ""))
    try:
        text = raw.decode("utf-8")
        files_out.append({"filename": safe_filename(os.path.basename(path), ""), "mimetype": "text/plain",
                           "data_base64": base64.b64encode(raw).decode()})
        preview = text if len(text) <= 4000 else text[:4000] + "\n...[truncated — full file attached below]"
        return f"Read `{path}` from {owner}/{repo} ({len(raw)} bytes):\n```\n{preview}\n```"
    except UnicodeDecodeError:
        return f"`{path}` in {owner}/{repo} isn't readable as text ({len(raw)} bytes, binary file)."


def do_github_commit_file(token, action_data, files_out):
    owner, repo, path = action_data["owner"], action_data["repo"], action_data["path"]
    content_b64 = base64.b64encode(action_data.get("content", "").encode("utf-8")).decode()
    sha = None
    try:
        existing = github_api(token, "GET", f"/repos/{owner}/{repo}/contents/{path}")
        sha = existing.get("sha")
    except Exception:
        pass  # file doesn't exist yet — that's fine, this becomes a create
    body = {"message": action_data.get("message", "Devgent commit"), "content": content_b64,
            "branch": action_data.get("branch", "main")}
    if sha:
        body["sha"] = sha
    result = github_api(token, "PUT", f"/repos/{owner}/{repo}/contents/{path}", json=body)
    html_url = (result.get("content") or {}).get("html_url", "")
    return f"Committed `{path}` to {owner}/{repo}. {html_url}"


def do_github_delete_file(token, action_data, files_out):
    owner, repo, path = action_data["owner"], action_data["repo"], action_data["path"]
    existing = github_api(token, "GET", f"/repos/{owner}/{repo}/contents/{path}")
    body = {"message": action_data.get("message", f"Devgent: delete {path}"),
            "sha": existing["sha"], "branch": action_data.get("branch", "main")}
    github_api(token, "DELETE", f"/repos/{owner}/{repo}/contents/{path}", json=body)
    return f"Deleted `{path}` from {owner}/{repo}."


def do_github_create_branch(token, action_data, files_out):
    owner, repo, branch = action_data["owner"], action_data["repo"], action_data["branch"]
    base = action_data.get("base", "main")
    base_ref = github_api(token, "GET", f"/repos/{owner}/{repo}/git/ref/heads/{base}")
    sha = base_ref["object"]["sha"]
    github_api(token, "POST", f"/repos/{owner}/{repo}/git/refs",
               json={"ref": f"refs/heads/{branch}", "sha": sha})
    return f"Created branch `{branch}` from `{base}` in {owner}/{repo}."


def do_github_list_pull_requests(token, action_data, files_out):
    owner, repo = action_data["owner"], action_data["repo"]
    state = action_data.get("state", "open")
    prs = github_api(token, "GET", f"/repos/{owner}/{repo}/pulls?state={state}")
    lines = [f"- #{p['number']} {p['title']} ({p['head']['ref']} -> {p['base']['ref']})" for p in prs[:20]]
    return f"{state.capitalize()} pull requests in {owner}/{repo}:\n" + ("\n".join(lines) if lines else "(none)")


def do_github_create_pull_request(token, action_data, files_out):
    owner, repo = action_data["owner"], action_data["repo"]
    body = {"title": action_data["title"], "head": action_data["head"], "base": action_data["base"],
            "body": action_data.get("body", "")}
    result = github_api(token, "POST", f"/repos/{owner}/{repo}/pulls", json=body)
    return f"Opened PR #{result['number']} in {owner}/{repo}: {result.get('html_url', '')}"


def do_github_merge_pull_request(token, action_data, files_out):
    owner, repo, number = action_data["owner"], action_data["repo"], action_data["pr_number"]
    body = {"merge_method": action_data.get("merge_method", "merge")}
    github_api(token, "PUT", f"/repos/{owner}/{repo}/pulls/{number}/merge", json=body)
    return f"Merged PR #{number} in {owner}/{repo}."


def do_github_close_pull_request(token, action_data, files_out):
    owner, repo, number = action_data["owner"], action_data["repo"], action_data["pr_number"]
    github_api(token, "PATCH", f"/repos/{owner}/{repo}/pulls/{number}", json={"state": "closed"})
    return f"Closed PR #{number} in {owner}/{repo} (not merged)."


def do_github_list_issues(token, action_data, files_out):
    owner, repo = action_data["owner"], action_data["repo"]
    state = action_data.get("state", "open")
    issues = github_api(token, "GET", f"/repos/{owner}/{repo}/issues?state={state}")
    issues = [i for i in issues if "pull_request" not in i]
    lines = [f"- #{i['number']} {i['title']}" for i in issues[:20]]
    return f"{state.capitalize()} issues in {owner}/{repo}:\n" + ("\n".join(lines) if lines else "(none)")


def do_github_create_issue(token, action_data, files_out):
    owner, repo = action_data["owner"], action_data["repo"]
    body = {"title": action_data["title"], "body": action_data.get("body", "")}
    result = github_api(token, "POST", f"/repos/{owner}/{repo}/issues", json=body)
    return f"Opened issue #{result['number']} in {owner}/{repo}: {result.get('html_url', '')}"


def do_github_comment_on_issue(token, action_data, files_out):
    owner, repo, number = action_data["owner"], action_data["repo"], action_data["issue_number"]
    github_api(token, "POST", f"/repos/{owner}/{repo}/issues/{number}/comments",
               json={"body": action_data.get("body", "")})
    return f"Commented on issue #{number} in {owner}/{repo}."


def do_github_close_issue(token, action_data, files_out):
    owner, repo, number = action_data["owner"], action_data["repo"], action_data["issue_number"]
    github_api(token, "PATCH", f"/repos/{owner}/{repo}/issues/{number}", json={"state": "closed"})
    return f"Closed issue #{number} in {owner}/{repo}."


def slack_api(token, method_name, **params):
    resp = requests.post(f"https://slack.com/api/{method_name}",
                          headers={"Authorization": f"Bearer {token}"}, json=params, timeout=30)
    data = resp.json()
    if not data.get("ok"):
        raise RuntimeError(f"Slack API error: {data.get('error', 'unknown_error')}")
    return data


def do_slack_list_channels(token, action_data, files_out):
    data = slack_api(token, "conversations.list", limit=50, types="public_channel,private_channel")
    names = ", ".join("#" + c["name"] for c in data.get("channels", [])[:30]) or "(no channels found)"
    return f"Channels I can see: {names}"


def do_slack_read_channel_messages(token, action_data, files_out):
    data = slack_api(token, "conversations.history", channel=action_data["channel"],
                      limit=action_data.get("limit", 20))
    lines = [f"- {m.get('user', 'someone')}: {m.get('text', '')[:200]}" for m in data.get("messages", [])]
    return f"Recent messages in {action_data['channel']}:\n" + ("\n".join(lines) if lines else "(none)")


def do_slack_send_message(token, action_data, files_out):
    data = slack_api(token, "chat.postMessage", channel=action_data["channel"], text=action_data.get("text", ""))
    return f"Sent to {data.get('channel', action_data['channel'])} in Slack."


def do_slack_reply_thread(token, action_data, files_out):
    data = slack_api(token, "chat.postMessage", channel=action_data["channel"],
                      thread_ts=action_data["thread_ts"], text=action_data.get("text", ""))
    return f"Replied in thread on {data.get('channel', action_data['channel'])}."


def do_slack_add_reaction(token, action_data, files_out):
    slack_api(token, "reactions.add", channel=action_data["channel"],
              timestamp=action_data["timestamp"], name=action_data.get("emoji", "thumbsup"))
    return f"Reacted with :{action_data.get('emoji', 'thumbsup')}: in {action_data['channel']}."


def notion_api(token, method, path, **kwargs):
    headers = {"Authorization": f"Bearer {token}", "Notion-Version": "2022-06-28"}
    resp = requests.request(method, f"https://api.notion.com/v1{path}", headers=headers, timeout=30, **kwargs)
    if resp.status_code >= 400:
        raise RuntimeError(f"Notion API error ({resp.status_code}): {resp.text[:300]}")
    return resp.json()


def do_notion_search(token, action_data, files_out):
    data = notion_api(token, "POST", "/search", json={"query": action_data.get("query", "")})
    results = data.get("results", [])[:10]
    lines = []
    for r in results:
        title_parts = (r.get("properties", {}).get("title", {}) or {}).get("title", [])
        title = "".join(p.get("plain_text", "") for p in title_parts) or "(untitled)"
        lines.append(f"- {title} (id: {r.get('id')})")
    return "Notion pages found:\n" + ("\n".join(lines) if lines else "(none)")


def do_notion_read_page(token, action_data, files_out):
    page_id = action_data["page_id"]
    blocks = notion_api(token, "GET", f"/blocks/{page_id}/children?page_size=50")
    lines = []
    for b in blocks.get("results", []):
        btype = b.get("type", "")
        rich_text = (b.get(btype, {}) or {}).get("rich_text", [])
        text = "".join(t.get("plain_text", "") for t in rich_text)
        if text:
            lines.append(text)
    content = "\n".join(lines) or "(no readable text content)"
    return f"Content of Notion page {page_id}:\n{content[:4000]}"


def do_notion_update_page(token, action_data, files_out):
    page_id = action_data["page_id"]
    children = [{"object": "block", "type": "paragraph",
                 "paragraph": {"rich_text": [{"type": "text", "text": {"content": p}}]}}
                for p in action_data.get("paragraphs", [])]
    notion_api(token, "PATCH", f"/blocks/{page_id}/children", json={"children": children})
    return f"Appended {len(children)} paragraph(s) to Notion page {page_id}."


def do_notion_create_page(token, action_data, files_out):
    parent = ({"page_id": action_data["parent_page_id"]} if action_data.get("parent_page_id")
              else {"workspace": True})
    children = [{"object": "block", "type": "paragraph",
                 "paragraph": {"rich_text": [{"type": "text", "text": {"content": p}}]}}
                for p in action_data.get("paragraphs", [])]
    body = {"parent": parent,
            "properties": {"title": [{"text": {"content": action_data.get("title", "Untitled")}}]},
            "children": children}
    result = notion_api(token, "POST", "/pages", json=body)
    return f"Created Notion page \"{action_data.get('title', 'Untitled')}\": {result.get('url', '')}"


def do_notion_archive_page(token, action_data, files_out):
    page_id = action_data["page_id"]
    notion_api(token, "PATCH", f"/pages/{page_id}", json={"archived": True})
    return f"Archived (deleted) Notion page {page_id}."


def figma_api(token, method, path, **kwargs):
    resp = requests.request(method, f"https://api.figma.com/v1{path}",
                             headers={"Authorization": f"Bearer {token}"}, timeout=30, **kwargs)
    if resp.status_code >= 400:
        raise RuntimeError(f"Figma API error ({resp.status_code}): {resp.text[:300]}")
    return resp.json()


def do_figma_get_file(token, action_data, files_out):
    data = figma_api(token, "GET", f"/files/{action_data['file_key']}")
    top_frames = [n.get("name") for n in (data.get("document", {}).get("children", []) or [])[:15]]
    return f"Figma file \"{data.get('name', '')}\" — top-level frames/pages: {', '.join(top_frames) or '(none found)'}"


def do_figma_list_comments(token, action_data, files_out):
    data = figma_api(token, "GET", f"/files/{action_data['file_key']}/comments")
    lines = [f"- {(c.get('user') or {}).get('handle', 'someone')}: {c.get('message', '')[:200]}"
             for c in data.get("comments", [])[:20]]
    return "Figma comments:\n" + ("\n".join(lines) if lines else "(none)")


def do_figma_add_comment(token, action_data, files_out):
    body = {"message": action_data.get("message", "")}
    if action_data.get("node_id"):
        body["client_meta"] = {"node_id": action_data["node_id"], "node_offset": {"x": 0, "y": 0}}
    figma_api(token, "POST", f"/files/{action_data['file_key']}/comments", json=body)
    return f"Added a comment to Figma file {action_data['file_key']}."


def do_figma_export_images(token, action_data, files_out):
    ids = ",".join(action_data.get("node_ids", []))
    fmt = action_data.get("format", "png")
    data = figma_api(token, "GET", f"/images/{action_data['file_key']}?ids={ids}&format={fmt}")
    images = data.get("images", {}) or {}
    count = 0
    for node_id, url in images.items():
        if not url:
            continue
        img_resp = requests.get(url, timeout=30)
        if img_resp.status_code == 200:
            fname = safe_filename(f"figma_{node_id.replace(':', '-')}", f".{fmt}")
            mimetype = "image/png" if fmt == "png" else ("image/svg+xml" if fmt == "svg" else "image/jpeg")
            files_out.append({"filename": fname, "mimetype": mimetype,
                               "data_base64": base64.b64encode(img_resp.content).decode()})
            count += 1
    return f"Exported {count} image(s) from Figma."


def ms_graph_api(token, method, path, **kwargs):
    resp = requests.request(method, f"https://graph.microsoft.com/v1.0{path}",
                             headers={"Authorization": f"Bearer {token}"}, timeout=30, **kwargs)
    if resp.status_code >= 400:
        raise RuntimeError(f"Microsoft Graph error ({resp.status_code}): {resp.text[:300]}")
    return resp.json() if resp.text else {}


def do_ms365_recent_files(token, action_data, files_out):
    data = ms_graph_api(token, "GET", "/me/drive/recent")
    names = [it.get("name") for it in data.get("value", [])[:15]]
    return "Recent OneDrive/SharePoint files: " + (", ".join(names) if names else "(none found)")


def do_ms365_search_files(token, action_data, files_out):
    from urllib.parse import quote
    q = quote(action_data.get("query", ""))
    data = ms_graph_api(token, "GET", f"/me/drive/root/search(q='{q}')")
    names = [it.get("name") for it in data.get("value", [])[:15]]
    return f"Files matching \"{action_data.get('query', '')}\": " + (", ".join(names) if names else "(none found)")


def do_ms365_read_file(token, action_data, files_out):
    item_id = action_data["item_id"]
    meta = ms_graph_api(token, "GET", f"/me/drive/items/{item_id}")
    resp = requests.get(f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content",
                         headers={"Authorization": f"Bearer {token}"}, timeout=30)
    if resp.status_code >= 400:
        raise RuntimeError(f"Microsoft Graph error ({resp.status_code}): {resp.text[:300]}")
    fname = meta.get("name", "file")
    try:
        text = resp.content.decode("utf-8")
        files_out.append({"filename": safe_filename(fname, ""), "mimetype": "text/plain",
                           "data_base64": base64.b64encode(resp.content).decode()})
        preview = text if len(text) <= 4000 else text[:4000] + "\n...[truncated — full file attached below]"
        return f"Read `{fname}` ({len(resp.content)} bytes):\n```\n{preview}\n```"
    except UnicodeDecodeError:
        files_out.append({"filename": safe_filename(fname, ""),
                           "mimetype": meta.get("file", {}).get("mimeType", "application/octet-stream"),
                           "data_base64": base64.b64encode(resp.content).decode()})
        return f"`{fname}` isn't plain text ({len(resp.content)} bytes) — attached as a file instead."


def do_ms365_list_teams_channels(token, action_data, files_out):
    if action_data.get("team_id"):
        data = ms_graph_api(token, "GET", f"/teams/{action_data['team_id']}/channels")
        names = [f"{c.get('displayName')} (id: {c.get('id')})" for c in data.get("value", [])]
        return f"Channels in team {action_data['team_id']}:\n" + ("\n".join(f"- {n}" for n in names) if names else "(none)")
    data = ms_graph_api(token, "GET", "/me/joinedTeams")
    names = [f"{t.get('displayName')} (id: {t.get('id')})" for t in data.get("value", [])]
    return "Teams I'm a member of:\n" + ("\n".join(f"- {n}" for n in names) if names else "(none found)")


def do_ms365_teams_messages(token, action_data, files_out):
    team_id, channel_id = action_data.get("team_id"), action_data.get("channel_id")
    data = ms_graph_api(token, "GET", f"/teams/{team_id}/channels/{channel_id}/messages")
    msgs = data.get("value", [])[:10]
    lines = []
    for m in msgs:
        body = (m.get("body", {}) or {}).get("content", "")
        author = ((m.get("from", {}) or {}).get("user", {}) or {}).get("displayName", "someone")
        lines.append(f"- {author}: {re.sub('<[^<]+?>', '', body)[:200]}")
    return "Recent Teams messages:\n" + ("\n".join(lines) if lines else "(none found)")


def do_ms365_send_teams_message(token, action_data, files_out):
    team_id, channel_id = action_data["team_id"], action_data["channel_id"]
    body = {"body": {"content": action_data.get("text", "")}}
    ms_graph_api(token, "POST", f"/teams/{team_id}/channels/{channel_id}/messages", json=body)
    return f"Posted a message to the Teams channel."


# Registry: action name -> (which CONNECTORS key it needs, handler(token, action_data, files_out) -> str)
CONNECTOR_ACTIONS = {
    "github_list_repos": ("github", do_github_list_repos),
    "github_list_files": ("github", do_github_list_files),
    "github_read_file": ("github", do_github_read_file),
    "github_commit_file": ("github", do_github_commit_file),
    "github_delete_file": ("github", do_github_delete_file),
    "github_create_branch": ("github", do_github_create_branch),
    "github_list_pull_requests": ("github", do_github_list_pull_requests),
    "github_create_pull_request": ("github", do_github_create_pull_request),
    "github_merge_pull_request": ("github", do_github_merge_pull_request),
    "github_close_pull_request": ("github", do_github_close_pull_request),
    "github_list_issues": ("github", do_github_list_issues),
    "github_create_issue": ("github", do_github_create_issue),
    "github_comment_on_issue": ("github", do_github_comment_on_issue),
    "github_close_issue": ("github", do_github_close_issue),

    "slack_list_channels": ("slack", do_slack_list_channels),
    "slack_read_channel_messages": ("slack", do_slack_read_channel_messages),
    "slack_send_message": ("slack", do_slack_send_message),
    "slack_reply_thread": ("slack", do_slack_reply_thread),
    "slack_add_reaction": ("slack", do_slack_add_reaction),

    "notion_search": ("notion", do_notion_search),
    "notion_read_page": ("notion", do_notion_read_page),
    "notion_create_page": ("notion", do_notion_create_page),
    "notion_update_page": ("notion", do_notion_update_page),
    "notion_archive_page": ("notion", do_notion_archive_page),

    "figma_get_file": ("figma", do_figma_get_file),
    "figma_list_comments": ("figma", do_figma_list_comments),
    "figma_add_comment": ("figma", do_figma_add_comment),
    "figma_export_images": ("figma", do_figma_export_images),

    "ms365_recent_files": ("microsoft365", do_ms365_recent_files),
    "ms365_search_files": ("microsoft365", do_ms365_search_files),
    "ms365_read_file": ("microsoft365", do_ms365_read_file),
    "ms365_list_teams_channels": ("microsoft365", do_ms365_list_teams_channels),
    "ms365_teams_messages": ("microsoft365", do_ms365_teams_messages),
    "ms365_send_teams_message": ("microsoft365", do_ms365_send_teams_message),
}


# =============================================================================
# Action execution — turns the model's JSON block into downloadable files
# =============================================================================

def execute_actions(actions, memory, connections=None):
    files = []
    system_notes = []
    connector_suggestions = []

    for item in actions:
        action_type = item.get("action")
        action_data = item.get("action_data", {}) or {}

        try:
            if action_type == "manage_memory":
                for k in action_data.get("discards", []):
                    memory.pop(k, None)
                for k, v in (action_data.get("updates", {}) or {}).items():
                    memory[k] = v

            elif action_type == "create_txt":
                fname = safe_filename(action_data.get("filename", "output.txt"), "")
                files.append({"filename": fname, "mimetype": "text/plain",
                              "data_base64": base64.b64encode(action_data.get("content", "").encode("utf-8")).decode()})

            elif action_type == "create_docx":
                fname = safe_filename(action_data.get("filename", "Document.docx"), ".docx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                              "data_base64": base64.b64encode(build_docx_bytes(action_data)).decode()})

            elif action_type == "create_pdf":
                fname = safe_filename(action_data.get("filename", "Document.pdf"), ".pdf")
                files.append({"filename": fname, "mimetype": "application/pdf",
                              "data_base64": base64.b64encode(build_pdf_bytes(action_data)).decode()})

            elif action_type == "create_xlsx":
                fname = safe_filename(action_data.get("filename", "Spreadsheet.xlsx"), ".xlsx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                              "data_base64": base64.b64encode(build_xlsx_bytes(action_data)).decode()})

            elif action_type == "create_pptx":
                fname = safe_filename(action_data.get("filename", "Presentation.pptx"), ".pptx")
                files.append({"filename": fname,
                              "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                              "data_base64": base64.b64encode(build_pptx_bytes(action_data)).decode()})

            elif action_type == "create_batch_files":
                zip_name = safe_filename(action_data.get("zip_filename", "devgent_files.zip"), ".zip")
                batch_files = action_data.get("files", [])
                files.append({"filename": zip_name, "mimetype": "application/zip",
                              "data_base64": base64.b64encode(build_batch_zip_bytes(batch_files)).decode()})
                # Also surface each file individually (same "group" = the zip name) so the chat UI
                # can copy/preview them on their own, and so an .html file's live preview can pull in
                # sibling .js/.css content from this same batch instead of only offering the .zip.
                for f in batch_files:
                    fname = safe_filename(f.get("filename", "file.txt"), "")
                    files.append({"filename": fname, "mimetype": "text/plain",
                                  "data_base64": base64.b64encode(f.get("content", "").encode("utf-8")).decode(),
                                  "group": zip_name})

            elif action_type in ("run_command", "start_dev_server"):
                system_notes.append(
                    f"⚠️ The `{action_type}` action isn't available in the web version for security reasons "
                    f"(a public server can't safely run arbitrary commands). Download the Devgent desktop app "
                    f"for terminal and dev-server support."
                )

            elif action_type == "suggest_connector":
                tool = action_data.get("tool")
                if tool in CONNECTORS:
                    connector_suggestions.append({"tool": tool, "reason": action_data.get("reason", "")})

            elif action_type in CONNECTOR_ACTIONS:
                tool, handler_fn = CONNECTOR_ACTIONS[action_type]
                token = _connector_token(connections, tool)
                system_notes.append(f"🔗 {handler_fn(token, action_data, files)}")

            else:
                system_notes.append(f"⚠️ Unknown action `{action_type}` — ignored.")

        except Exception as e:
            system_notes.append(f"❌ Failed to complete `{action_type}`: {e}")

    return files, memory, system_notes, connector_suggestions


# =============================================================================
# HTTP handler (Vercel Python runtime contract)
# =============================================================================

class handler(BaseHTTPRequestHandler):

    def _cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def do_OPTIONS(self):
        self.send_response(204)
        self._cors_headers()
        self.end_headers()

    def _base_url(self):
        proto = self.headers.get("x-forwarded-proto", "https")
        host = self.headers.get("x-forwarded-host") or self.headers.get("Host", "")
        path = self.path.split("?", 1)[0]
        return f"{proto}://{host}{path}"

    def _query(self):
        from urllib.parse import urlparse, parse_qs
        return {k: v[0] for k, v in parse_qs(urlparse(self.path).query).items()}

    def do_GET(self):
        q = self._query()
        base_url = self._base_url()

        # Step 1: kick the browser into the provider's own OAuth consent screen.
        if q.get("connect") in CONNECTORS:
            tool = q["connect"]
            url = oauth_authorize_url(tool, base_url, state=q.get("state", ""))
            if not url:
                return self._respond(200, {
                    "error": f"{CONNECTORS[tool]['label']} isn't configured on this deployment yet.",
                    "setup_help": CONNECTORS[tool]["setup_help"],
                })
            self.send_response(302)
            self._cors_headers()
            self.send_header("Location", url)
            self.end_headers()
            return

        # Step 2: provider redirects back here with a ?code=; exchange it server-side and
        # hand the resulting access token to the popup's opener via postMessage.
        if q.get("oauth_callback") in CONNECTORS:
            tool = q["oauth_callback"]
            try:
                if q.get("error"):
                    raise RuntimeError(q.get("error_description", q["error"]))
                payload = oauth_exchange_code(tool, q.get("code", ""), base_url)
                message = {"devgent_oauth": True, "ok": True, "tool": tool,
                           "access_token": payload.get("access_token"),
                           "team": (payload.get("team") or {}).get("name") if tool == "slack" else None}
                body_note = f"Connected {CONNECTORS[tool]['label']} — you can close this tab."
            except Exception as e:
                message = {"devgent_oauth": True, "ok": False, "tool": tool, "error": str(e)}
                body_note = f"Couldn't connect {CONNECTORS[tool]['label']}: {e}"
            html = f"""<!DOCTYPE html><html><body style="font-family:sans-serif;background:#0D1420;color:#F1F5F9;
display:flex;align-items:center;justify-content:center;height:100vh;margin:0;">
<p>{body_note}</p>
<script>
  if (window.opener) {{ window.opener.postMessage({json.dumps(message)}, "*"); }}
  setTimeout(function(){{ window.close(); }}, 800);
</script>
</body></html>"""
            self.send_response(200)
            self._cors_headers()
            self.send_header("Content-Type", "text/html")
            self.end_headers()
            self.wfile.write(html.encode("utf-8"))
            return

        self.send_response(200)
        self._cors_headers()
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps({"status": "Devgent backend is online"}).encode())

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            raw_body = self.rfile.read(length) if length else b"{}"
            payload = json.loads(raw_body or b"{}")

            message = (payload.get("message") or "").strip()
            history = payload.get("history") or []
            api_key = (payload.get("apiKey") or "").strip()
            model = payload.get("model") or "gemini-3.6-flash"
            memory = payload.get("memory") or {}
            attachment = payload.get("attachment")
            connections = payload.get("connections") or {}

            if not api_key:
                if message.strip().lower() == "help":
                    reply = GEMINI_API_KEY_HELP_TEXT
                else:
                    reply = "Please enter a free Gemini API key at the top of the page. Type 'help' if you want to know how to get one."
                return self._respond(200, {"reply": reply, "files": [], "memory": memory,
                                            "connector_suggestions": [], "connectors": connected_tools(connections)})

            if not message and not attachment:
                return self._respond(400, {"error": "Empty message."})

            full_prompt = message
            if attachment and attachment.get("data_base64"):
                try:
                    raw_bytes = base64.b64decode(attachment["data_base64"])
                    extracted = extract_file_content(attachment.get("filename", "attachment"), raw_bytes)
                    if extracted:
                        full_prompt += f"\n\n[Attached File Content ({attachment.get('filename', 'attachment')})]:\n{extracted[:20000]}"
                except Exception as e:
                    full_prompt += f"\n\n[Could not read attached file: {e}]"

            full_prompt += (
                "\n\n[System reminder: If you learned anything new about the user, the project, "
                "their preferences, or task status this turn, include a manage_memory action to save it - "
                "even if the user did not explicitly ask you to remember it.]"
            )

            memory_str = json.dumps(memory, indent=2) if memory else "No stored memories yet."
            connected = connected_tools(connections)
            connectors_lines = [
                f"- {cfg['label']} ({key}): {'CONNECTED' if key in connected else 'not connected'}"
                for key, cfg in CONNECTORS.items()
            ]
            sys_inst = SYSTEM_INSTRUCTION_TEMPLATE.format(
                current_date=date.today().isoformat(),
                memory_state=memory_str,
                connectors_state="\n".join(connectors_lines),
            )

            response_text = call_gemini(api_key, model, sys_inst, history, full_prompt)
            response_text = re.sub(r'<think>[\s\S]*?</think>', '', response_text, flags=re.IGNORECASE).strip()

            # Pull the trailing ```json ... ``` action block out of the reply, same approach as the desktop app.
            json_match = re.search(r"```json\s*(\{[\s\S]*?\})\s*```", response_text)
            reply_text = response_text
            actions = []
            if json_match:
                reply_text = response_text[:json_match.start()].strip()
                try:
                    parsed = json.loads(json_match.group(1))
                    actions = parsed.get("actions", [])
                except Exception:
                    pass

            files, memory, system_notes, connector_suggestions = (
                execute_actions(actions, memory, connections) if actions else ([], memory, [], [])
            )
            if system_notes:
                reply_text = (reply_text + "\n\n" + "\n".join(system_notes)).strip()

            self._respond(200, {"reply": reply_text, "files": files, "memory": memory,
                                 "connector_suggestions": connector_suggestions, "connectors": connected})

        except Exception as e:
            traceback.print_exc()
            self._respond(500, {"error": str(e)})

    def _respond(self, status, obj):
        self.send_response(status)
        self._cors_headers()
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps(obj).encode("utf-8"))
