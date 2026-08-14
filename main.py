import os
import sys
import json
import re
import requests
import threading
import subprocess
import socket
import webbrowser
import http.server
import socketserver
import platform
import traceback
from pathlib import Path
# Modern GUI Framework (PyQt6)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize
from PyQt6.QtGui import QFont, QIcon, QPixmap, QPainter, QBrush, QColor
from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QPushButton, QTextEdit, QTextBrowser, QLineEdit, QComboBox, QFrame,
    QScrollArea, QFileDialog, QMessageBox
)

# Markdown rendering for the chat window (renders headings, bold/italic, lists,
# tables, and fenced ```code``` blocks as real HTML instead of raw text).
try:
    import markdown2
    MARKDOWN2_AVAILABLE = True
except ImportError:
    MARKDOWN2_AVAILABLE = False

# Document & Spreadsheet Extractors / Generators, and the Gemini API client
# are intentionally NOT imported here at module load time. Importing
# python-docx, openpyxl, pypdf, reportlab, python-pptx, and (especially)
# google-genai eagerly used to add 10-20+ seconds before the main window
# ever appeared - all of that work happened before QApplication was even
# created. They're now imported lazily, on first actual use, via
# _ensure_heavy_imports() below, and warmed up in a background thread at
# startup so the window shows almost instantly. See _ensure_heavy_imports().

# ==========================================
# CONFIGURATION (BYOK - Bring Your Own Key)
# ==========================================
DEVGENT_COMPLETE = "no"

# Devgent no longer ships with any hardcoded API key. Each user supplies their
# own Google Gemini API key at first run (or via the GEMINI_API_KEY environment
# variable, which always takes priority). The key is persisted locally to a
# small config file in the user's home directory so it survives app restarts -
# this is intentionally OUTSIDE the PyInstaller-built executable/app bundle so
# rebuilding or redistributing the .exe never leaks anyone's key.
CONFIG_DIR = os.path.join(str(Path.home()), ".devgent")
CONFIG_FILE = os.path.join(CONFIG_DIR, "config.json")


def load_saved_api_key():
    """Loads the Gemini API key: env var takes priority, then the local config file."""
    env_key = os.environ.get("GEMINI_API_KEY", "").strip()
    if env_key:
        return env_key
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                return str(data.get("gemini_api_key", "")).strip()
        except Exception:
            return ""
    return ""


def save_api_key_to_disk(key):
    """Persists the Gemini API key locally so the user only has to enter it once."""
    try:
        os.makedirs(CONFIG_DIR, exist_ok=True)
        existing = {}
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                    existing = json.load(f)
            except Exception:
                existing = {}
        existing["gemini_api_key"] = key
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(existing, f, indent=2)
        return True
    except Exception as e:
        print(f"[CONSOLE DEBUG] Failed to save API key: {e}")
        return False


_heavy_imports_lock = threading.Lock()
_heavy_imports_done = False


def _ensure_heavy_imports():
    """Imports python-docx, openpyxl, pypdf, reportlab, python-pptx, and the
    Gemini API client on first actual use, instead of at module load time.

    This is what makes the window appear instantly instead of 10-20 seconds
    after launch: those libraries (google-genai especially) are slow to
    import, and previously that cost was paid before QApplication was even
    constructed. Now it's paid the first time a document is generated, a
    file is attached, or a chat message is sent - and main() also fires this
    off in a background thread as soon as the app starts, so by the time the
    user actually needs one of these features it has usually already
    finished loading. Safe to call from any thread, and safe to call more
    than once (it's a no-op after the first successful run).
    """
    global _heavy_imports_done
    global docx, DocxPt, DocxRGBColor, DocxInches, DocxEmu, OxmlElement, qn
    global openpyxl, Font, PatternFill, Alignment, Border, Side, get_column_letter
    global PdfReader
    global letter, inch, colors, getSampleStyleSheet, ParagraphStyle
    global SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    global Presentation, Inches, Pt, MSO_SHAPE, PP_ALIGN
    global XL_CHART_TYPE, XL_LEGEND_POSITION, CategoryChartData, RGBColor
    global genai, types

    if _heavy_imports_done:
        return

    with _heavy_imports_lock:
        if _heavy_imports_done:
            return

        import docx as _docx
        from docx.shared import Pt as _DocxPt, RGBColor as _DocxRGBColor, Inches as _DocxInches, Emu as _DocxEmu
        from docx.oxml import OxmlElement as _OxmlElement
        from docx.oxml.ns import qn as _qn

        import openpyxl as _openpyxl
        from openpyxl.styles import Font as _Font, PatternFill as _PatternFill, Alignment as _Alignment, Border as _Border, Side as _Side
        from openpyxl.utils import get_column_letter as _get_column_letter

        from pypdf import PdfReader as _PdfReader

        from reportlab.lib.pagesizes import letter as _letter
        from reportlab.lib.units import inch as _inch
        from reportlab.lib import colors as _colors
        from reportlab.lib.styles import getSampleStyleSheet as _getSampleStyleSheet, ParagraphStyle as _ParagraphStyle
        from reportlab.platypus import (
            SimpleDocTemplate as _SimpleDocTemplate, Paragraph as _Paragraph, Spacer as _Spacer,
            Table as _Table, TableStyle as _TableStyle, ListFlowable as _ListFlowable, ListItem as _ListItem
        )

        from pptx import Presentation as _Presentation
        from pptx.util import Inches as _Inches, Pt as _Pt
        from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN
        from pptx.enum.chart import XL_CHART_TYPE as _XL_CHART_TYPE, XL_LEGEND_POSITION as _XL_LEGEND_POSITION
        from pptx.chart.data import CategoryChartData as _CategoryChartData
        from pptx.dml.color import RGBColor as _RGBColor

        from google import genai as _genai
        from google.genai import types as _types

        docx, DocxPt, DocxRGBColor, DocxInches, DocxEmu, OxmlElement, qn = \
            _docx, _DocxPt, _DocxRGBColor, _DocxInches, _DocxEmu, _OxmlElement, _qn
        openpyxl, Font, PatternFill, Alignment, Border, Side, get_column_letter = \
            _openpyxl, _Font, _PatternFill, _Alignment, _Border, _Side, _get_column_letter
        PdfReader = _PdfReader
        letter, inch, colors, getSampleStyleSheet, ParagraphStyle = \
            _letter, _inch, _colors, _getSampleStyleSheet, _ParagraphStyle
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem = \
            _SimpleDocTemplate, _Paragraph, _Spacer, _Table, _TableStyle, _ListFlowable, _ListItem
        Presentation, Inches, Pt, MSO_SHAPE, PP_ALIGN = \
            _Presentation, _Inches, _Pt, _MSO_SHAPE, _PP_ALIGN
        XL_CHART_TYPE, XL_LEGEND_POSITION, CategoryChartData, RGBColor = \
            _XL_CHART_TYPE, _XL_LEGEND_POSITION, _CategoryChartData, _RGBColor
        genai, types = _genai, _types

        _heavy_imports_done = True


GEMINI_API_KEY_HELP_TEXT = """Getting a Gemini API key is completely free and only takes a minute:

1. Open Google AI Studio in your browser at https://aistudio.google.com/apikey
2. Sign in with your Google account if you're asked to.
3. Click "Create API key", then choose to create a new project (or pick an existing one) when prompted.
4. Copy the API key that gets generated.
5. Paste it into the "Gemini API Key" field in the sidebar on the left.

No credit card or payment is required - the free tier is enough to chat with Devgent."""

# System Instructions with Dynamic Memory Handling
SYSTEM_INSTRUCTION_TEMPLATE = """
You are Devgent, an AI desktop assistant executing tasks on {current_os} using {shell_name}.
Current Year: 2026

=== CORE OPERATING PRINCIPLES ===
1. Respond clearly to the user in plain chat text, then, ONLY if you are actually taking an action this turn, attach a single JSON block at the very end of your response inside triple backticks tagged with json.
2. The JSON block can contain an "actions" list so you can execute multiple actions sequentially in a single turn.
3. File Naming: Always specify exact file extensions (.py, .js, .html, .css, .json, .txt, .pdf, .docx, .xlsx, .pptx).
4. Operating System Awareness: You are on {current_os}. When formulating terminal commands, use native {current_os} shell commands.
5. Permission Prompts: Every action that creates/overwrites a file, runs a terminal command, or launches the dev server will automatically show the user a Yes/No confirmation dialog before it happens. You do NOT need to separately beg permission in your own text - just briefly describe what you're about to propose, then issue the action. If the user clicks "No", you'll see that in the next system message and should not silently retry the same thing.
6. Markdown In Chat: Your plain-text reply (everything outside the JSON block) is rendered as real Markdown in the chat window - headings, **bold**, *italics*, bullet/numbered lists, tables, and fenced code blocks with a language tag (e.g. ```python ... ```) will all display properly formatted, including syntax-friendly code blocks. Use Markdown freely to make explanations, code snippets, and instructions easy to read. This is completely separate from the JSON action block used to trigger file/command actions.

=== CLARIFY BEFORE ACTING (VERY IMPORTANT) ===
If the user's request is vague, incomplete, or does not tell you concretely WHAT to build, write, or do - for example "I want to make a python script", "build me something cool", "make a presentation", "help me with a document" - then you must NOT take any action and must NOT emit a JSON action block at all this turn. Instead:
- Ask the user, in plain chat text only, exactly what they want: purpose, subject matter, desired content/sections, style, filename, etc. - whatever is missing.
- If they already gave you some partial detail, save it to memory with "manage_memory" so it isn't lost (this is the one exception where you may still emit a JSON block even without acting on a file/command), e.g.:
```json
{{
  "actions": [
    {{"action": "manage_memory", "action_data": {{"updates": {{"pending_request": "User wants a Python script, no further details yet."}}}}}}
  ]
}}
```
- Wait for the user's next message before creating, writing, or running anything.
Only move on to create_txt / create_docx / create_pdf / create_pptx / create_xlsx / create_batch_files / run_command / start_dev_server once the user has given you enough concrete specifics to act on.

=== OPERATING SYSTEM COMMAND TIPS ===
- To open a webpage in the default web browser:
  * Windows: Use the command `start <url>` (e.g., `start https://google.com`).
  * macOS: Use `open <url>` (e.g., `open https://google.com`).
  * Linux: Use `xdg-open <url>` (e.g., `xdg-open https://google.com`).
- To launch local files or applications:
  * Windows: `start <file_path>`
  * macOS: `open <file_path>`
  * Linux: `xdg-open <file_path>`

=== DYNAMIC MEMORY STORE (BE AGGRESSIVE ABOUT USING THIS) ===
Current Stored Memories:
{memory_state}

You can save, update, or discard memory items dynamically to retain user context across conversations. Treat this as a first-class habit, not an afterthought - on almost every turn, ask yourself "did I just learn something worth remembering?" That includes, but is not limited to:
- The user's name, role, or how they prefer to be addressed.
- Project name, purpose, tech stack, and key file/folder names already in use.
- Stylistic or technical preferences (indentation, frameworks, color themes, tone).
- Decisions the user has made or corrections they've given you (these should overwrite/replace old memory values, not just add new ones).
- The current status of a multi-step task, so you can resume correctly next session.
If the answer is yes, include a "manage_memory" action - it can run alongside any other action in the same "actions" array in the same turn, it does not need its own turn. Err on the side of saving too much rather than too little; unneeded or outdated memories can always be cleaned up later with "discards".

=== MULTI-ACTION SCHEMAS ===
You can run multiple actions in sequence by passing them in an "actions" array:
\\`\\`\\`json
{{
  "actions": [
    {{
      "action": "create_txt",
      "action_data": {{"filename": "app.py", "content": "print('hello')"}}
    }},
    {{
      "action": "run_command",
      "action_data": {{"command": "python app.py", "description": "Running app.py"}}
    }}
  ]
}}
\\`\\`\\`

=== AVAILABLE ACTIONS ===

1. Manage Dynamic Memory (manage_memory)
{{"action": "manage_memory", "action_data": {{"updates": {{"user_theme": "dark"}}, "discards": ["old_project_key"]}}}}
Remember that the updates: part of the manage_memory action checks if the memory category exists, and if it doesn't, you can create a new category with updates:

2. Advanced Styled PowerPoint Presentation (create_pptx)
Follow this schema and syntax exactly:
```json
{{
  "action": "create_pptx",
  "action_data": {{
    "filename": "pitch.pptx",
    "title": "Deck Title",
    "subtitle": "Optional subtitle line",
    "title_color": "#3B82F6",
    "subtitle_color": "#94A3B8",
    "background_color": "#0F172A",
    "theme_accent": "#3B82F6",
    "slides": [
      {{
        "layout": "bullets",
        "title": "Slide Heading",
        "background_color": "#F8FAFC",
        "header_color": "#1E293B",
        "header_text_color": "#FFFFFF",
        "accent_color": "#3B82F6",
        "content": ["First point", "Second point", "Third point"],
        "chart": {{
          "type": "column",
          "title": "Optional chart title",
          "categories": ["Q1", "Q2", "Q3"],
          "series": [{{"name": "Revenue", "values": [10, 20, 30]}}]
        }}
      }},
      {{
        "layout": "two_column",
        "title": "Comparison",
        "left_heading": "Pros",
        "content_left": ["Fast", "Cheap"],
        "right_heading": "Cons",
        "content_right": ["Riskier"]
      }},
      {{
        "layout": "quote",
        "title": "Customer Voice",
        "quote_text": "This product changed how we work.",
        "quote_author": "Jane Doe, CTO"
      }},
      {{
        "layout": "section",
        "title": "Part Two: Results",
        "subtitle": "Optional divider subtitle",
        "background_color": "#3B82F6",
        "header_text_color": "#FFFFFF"
      }}
    ]
  }}
}}
```
Rules:
- "layout" is one of: "bullets" (default; supports an optional "chart" side-by-side), "two_column", "quote", or "section" (a full-bleed divider slide with no body card).
- All colors are hex strings like "#0F172A". Any slide can override background_color, header_color, header_text_color, and accent_color; otherwise it inherits "theme_accent".
- Chart "type" is one of "bar", "column", "pie", "line".
- Every content slide automatically gets accent-colored bullet markers, a card background, and a footer with a page counter and a thin accent rule - you do not need to add these yourself.
- Prefer 3-6 short bullets per slide rather than long paragraphs.

3. Excel Spreadsheet (create_xlsx)
{{"action": "create_xlsx", "action_data": {{"filename": "data.xlsx", "sheet_name": "Metrics", "headers": ["Name", "Score"], "rows": [["Alice", 95]]}}}}

4. Word Document (create_docx)
{{"action": "create_docx", "action_data": {{"filename": "doc.docx", "title": "Proposal", "heading": "Summary", "paragraphs": ["Details..."]}}}}

5. Styled PDF Document (create_pdf)
Follow this schema and syntax exactly:
```json
{{
  "action": "create_pdf",
  "action_data": {{
    "filename": "report.pdf",
    "title": "Report Title",
    "subtitle": "Optional subtitle",
    "accent_color": "#3B82F6",
    "sections": [
      {{
        "heading": "Overview",
        "paragraphs": ["Body paragraph text goes here."],
        "bullets": ["Point one", "Point two"],
        "quote": "An optional pull-quote highlighting something important.",
        "table": {{"headers": ["Name", "Score"], "rows": [["Alice", "95"], ["Bob", "88"]]}}
      }}
    ]
  }}
}}
```
Rules:
- The title/subtitle render inside a colored banner using "accent_color".
- Each section's "heading" gets a colored underline; "bullets" render as a colored bullet list; "quote" renders as an indented, colored pull-quote; "table" (optional) renders as a striped, styled data table.
- For very simple documents you may skip "sections" entirely and just pass a flat "paragraphs" list, same as before - it will still get the styled title banner and footer.
- Every page automatically gets a footer with an accent-colored rule and a page number.

6. Single Text/Code File (create_txt)
{{"action": "create_txt", "action_data": {{"filename": "app.py", "content": "print('hello')"}}}}

7. Batch Web Files (create_batch_files)
{{"action": "create_batch_files", "action_data": {{"files": [{{"filename": "index.html", "content": "<h1>Hello</h1>"}}]}}}}

8. Run Terminal Command (run_command)
{{"action": "run_command", "action_data": {{"command": "start https://google.com", "description": "Opening Google in default browser"}}}}
This always requires the user's explicit Yes/No confirmation before it runs.

9. Launch Local Web Server (start_dev_server)
{{"action": "start_dev_server", "action_data": {{}}}}
This also requires the user's explicit Yes/No confirmation before it binds a port and opens a browser tab.

Remember to USE your memory as much as you can. Store context, the user's task, their name, everything you think you might want to use later. You will remember absolutely nothing that you do not store in your memory. It is recommended to store the context of your previous responses in your memory, too.
"""


def hex_to_rgb(hex_str):
    """Converts hex color string to RGBColor object."""
    hex_str = hex_str.lstrip('#')
    if len(hex_str) != 6:
        hex_str = "000000"
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def shade_color(hex_str, amount):
    """Lightens (amount > 0) or darkens (amount < 0) a hex color string.
    amount is a float in roughly [-1, 1]."""
    hex_str = hex_str.lstrip('#')
    if len(hex_str) != 6:
        hex_str = "3B82F6"
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    if amount >= 0:
        r = int(r + (255 - r) * amount)
        g = int(g + (255 - g) * amount)
        b = int(b + (255 - b) * amount)
    else:
        r = int(r * (1 + amount))
        g = int(g * (1 + amount))
        b = int(b * (1 + amount))
    r, g, b = max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b))
    return f"#{r:02X}{g:02X}{b:02X}"


def extract_file_content(file_path):
    """Extracts clean text from docx, xlsx, pdf, or text files."""
    if not file_path or not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.docx', '.xlsx', '.pdf'):
        _ensure_heavy_imports()

    try:
        if ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(file_path, data_only=True)
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        content.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            return "\n".join(content)

        elif ext == '.pdf':
            reader = PdfReader(file_path)
            return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

    except Exception as e:
        print(f"[CONSOLE DEBUG] File Extraction Error: {e}")
        sys.stdout.flush()
        return f"[Error extracting file content from {os.path.basename(file_path)}: {e}]"


class ReusableTCPServer(socketserver.TCPServer):
    allow_reuse_address = True


class CustomHTTPRequestHandler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *args, directory=None, **kwargs):
        super().__init__(*args, directory=directory, **kwargs)

    def do_GET(self):
        if self.path == "/" or not self.path:
            self.path = "/index.html"
        return super().do_GET()

    def log_message(self, format, *args):
        pass


class WorkerThread(QThread):
    finished_signal = pyqtSignal(str, object)
    error_signal = pyqtSignal(str)

    def __init__(self, task_func, *args, **kwargs):
        super().__init__()
        self.task_func = task_func
        self.args = args
        self.kwargs = kwargs

    def run(self):
        print("\n[CONSOLE DEBUG] WorkerThread: Starting thread execution...")
        sys.stdout.flush()
        try:
            text_res, payload = self.task_func(*self.args, **self.kwargs)
            print("[CONSOLE DEBUG] WorkerThread: Execution finished cleanly.")
            sys.stdout.flush()
            self.finished_signal.emit(text_res, payload)
        except Exception as e:
            print("\n[CONSOLE DEBUG - FATAL THREAD CRASH]")
            traceback.print_exc()
            sys.stdout.flush()
            self.error_signal.emit(str(e))


def generate_app_icon():
    pixmap = QPixmap(64, 64)
    pixmap.fill(Qt.GlobalColor.transparent)
    painter = QPainter(pixmap)
    painter.setRenderHint(QPainter.RenderHint.Antialiasing)

    painter.setBrush(QBrush(QColor("#3B82F6")))
    painter.setPen(Qt.PenStyle.NoPen)
    painter.drawRoundedRect(0, 0, 64, 64, 16, 16)

    painter.setBrush(QBrush(QColor("#0F172A")))
    painter.drawRoundedRect(12, 12, 40, 40, 10, 10)

    font = QFont("Segoe UI", 22, QFont.Weight.Bold)
    painter.setFont(font)
    painter.setPen(QColor("#60A5FA"))
    painter.drawText(0, 0, 64, 64, Qt.AlignmentFlag.AlignCenter, "</>")

    painter.end()
    return QIcon(pixmap)


class DevgentApp(QMainWindow):
    def __init__(self):
        super().__init__()

        self.selected_file_path = None
        self.output_directory = os.path.normpath(str(Path.home() / "Documents"))
        os.makedirs(self.output_directory, exist_ok=True)

        self.memory_file_path = os.path.join(self.output_directory, ".devgent_memory.json")
        self.agent_memory = self._load_memory()
        self.chat_history = []

        self.active_server_httpd = None
        self.active_server_port = None
        self.chat_bubbles = {}

        self.thinking_bubble_frame = None
        self.thinking_label = None

        self.active_workers = set()

        self.current_os = platform.system()
        self.shell_name = "Command Prompt / PowerShell" if self.current_os == "Windows" else "zsh / bash"

        self.gemini_api_key = load_saved_api_key()

        self.setWindowTitle("Devgent Desktop Agent (PyQt6 Edition)")
        self.setWindowIcon(generate_app_icon())
        self.resize(1280, 850)

        self._apply_stylesheet()
        self._build_ui()
        self.append_chat("Devgent", "Welcome back! Dynamic agent memory is online. What are we building today?")

    def _load_memory(self):
        if os.path.exists(self.memory_file_path):
            try:
                with open(self.memory_file_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                return {}
        return {}

    def _save_memory(self):
        try:
            with open(self.memory_file_path, "w", encoding="utf-8") as f:
                json.dump(self.agent_memory, f, indent=2)
        except Exception as e:
            print(f"[CONSOLE DEBUG] Error saving memory: {e}")

    def _get_system_instructions(self):
        mem_str = json.dumps(self.agent_memory, indent=2) if self.agent_memory else "No stored memories yet."
        return SYSTEM_INSTRUCTION_TEMPLATE.format(
            current_os=self.current_os,
            shell_name=self.shell_name,
            memory_state=mem_str
        ).replace("\\`\\`\\`", "```")

    def _apply_stylesheet(self):
        self.setStyleSheet("""
            QMainWindow { background-color: #0F172A; }
            QWidget { color: #F8FAFC; font-family: 'Segoe UI', Inter, sans-serif; }
            QFrame#Sidebar { background-color: #1E293B; border-right: 1px solid #334155; }
            QFrame#ChatContainer { background-color: #0F172A; border: 1px solid #334155; border-radius: 12px; }
            QFrame#InputDock { background-color: #1E293B; border: 1px solid #334155; border-radius: 10px; }
            QLineEdit, QTextEdit { background-color: #0F172A; border: 1px solid #334155; border-radius: 6px; padding: 8px; color: #F8FAFC; }
            QLineEdit:focus, QTextEdit:focus { border: 1px solid #3B82F6; }
            QPushButton { background-color: #3B82F6; color: #FFFFFF; border: none; border-radius: 6px; padding: 8px 14px; font-weight: bold; }
            QPushButton:hover { background-color: #2563EB; }
            QPushButton#SecondaryBtn { background-color: #334155; color: #F8FAFC; }
            QPushButton#SecondaryBtn:hover { background-color: #475569; }
            QPushButton#SuccessBtn { background-color: #10B981; }
            QPushButton#SuccessBtn:hover { background-color: #059669; }
            QComboBox { background-color: #0F172A; border: 1px solid #334155; border-radius: 6px; padding: 6px; color: #F8FAFC; }
            QScrollArea { border: none; background-color: #0F172A; }
            QScrollArea > QWidget > QWidget { background-color: #0F172A; }
            QWidget#ChatScrollContent { background-color: #0F172A; }
            QScrollBar:vertical { border: none; background: #0F172A; width: 8px; margin: 0px; border-radius: 4px; }
            QScrollBar::handle:vertical { background: #334155; min-height: 30px; border-radius: 4px; }
            QScrollBar::handle:vertical:hover { background: #475569; }
            QScrollBar::handle:vertical:active { background: #3B82F6; }
            QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical { border: none; background: none; height: 0px; }
            QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical { background: none; }
            QFrame#UserBubble QLabel, QFrame#DevgentBubble QLabel { border: none; background: transparent; }
        """)

    def _build_ui(self):
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        main_layout = QHBoxLayout(main_widget)
        main_layout.setContentsMargins(12, 12, 12, 12)
        main_layout.setSpacing(12)

        self.sidebar = QFrame()
        self.sidebar.setObjectName("Sidebar")
        self.sidebar.setFixedWidth(280)
        sidebar_layout = QVBoxLayout(self.sidebar)
        sidebar_layout.setContentsMargins(16, 16, 16, 16)
        sidebar_layout.setSpacing(12)

        brand_layout = QHBoxLayout()
        brand_icon = QLabel()
        brand_icon.setPixmap(generate_app_icon().pixmap(32, 32))
        brand_layout.addWidget(brand_icon)

        logo = QLabel("Devgent")
        logo.setFont(QFont("Segoe UI", 18, QFont.Weight.Bold))
        logo.setStyleSheet("color: #3B82F6;")
        brand_layout.addWidget(logo)
        brand_layout.addStretch()

        sidebar_layout.addLayout(brand_layout)

        sub_logo = QLabel("Desktop Agent Workspace")
        sub_logo.setStyleSheet("color: #94A3B8; font-size: 11px;")
        sidebar_layout.addWidget(sub_logo)

        sidebar_layout.addSpacing(10)

        engine_lbl = QLabel("Execution Engine")
        engine_lbl.setFont(QFont("Segoe UI", 10, QFont.Weight.Bold))
        sidebar_layout.addWidget(engine_lbl)

        self.engine_option = QComboBox()
        self.engine_option.addItems([
            "Devgent Ultra Heavyweight (gemini-3.6-flash)",
            "Devgent Turbo Speed (gemini-3.5-flash-lite)",
            "Ollama (Local deepseek-coder)"
        ])
        sidebar_layout.addWidget(self.engine_option)

        sidebar_layout.addSpacing(10)

        key_lbl = QLabel("Gemini API Key")
        key_lbl.setFont(QFont("Segoe UI", 10, QFont.Weight.Bold))
        sidebar_layout.addWidget(key_lbl)

        self.api_key_input = QLineEdit()
        self.api_key_input.setEchoMode(QLineEdit.EchoMode.Password)
        self.api_key_input.setPlaceholderText("Paste your free Gemini API key")
        self.api_key_input.setText(self.gemini_api_key)
        self.api_key_input.editingFinished.connect(self._on_api_key_changed)
        sidebar_layout.addWidget(self.api_key_input)

        key_hint = QLabel("It's free \u2014 reply 'help' in chat if you need one.")
        key_hint.setWordWrap(True)
        key_hint.setStyleSheet("color: #94A3B8; font-size: 10px;")
        sidebar_layout.addWidget(key_hint)

        sidebar_layout.addSpacing(10)

        dir_lbl = QLabel("Workspace Location")
        dir_lbl.setFont(QFont("Segoe UI", 10, QFont.Weight.Bold))
        sidebar_layout.addWidget(dir_lbl)

        self.browse_btn = QPushButton("📁 Select Folder")
        self.browse_btn.setObjectName("SecondaryBtn")
        self.browse_btn.clicked.connect(self.select_output_directory)
        sidebar_layout.addWidget(self.browse_btn)

        self.path_display = QLabel(self.output_directory)
        self.path_display.setWordWrap(True)
        self.path_display.setStyleSheet("color: #94A3B8; font-size: 10px;")
        sidebar_layout.addWidget(self.path_display)

        sidebar_layout.addSpacing(10)

        self.server_btn = QPushButton("🌐 Launch Live Server")
        self.server_btn.setObjectName("SuccessBtn")
        self.server_btn.clicked.connect(self.start_dev_server)
        sidebar_layout.addWidget(self.server_btn)

        sidebar_layout.addSpacing(10)

        mem_lbl = QLabel("Active Persistent Memory")
        mem_lbl.setFont(QFont("Segoe UI", 10, QFont.Weight.Bold))
        sidebar_layout.addWidget(mem_lbl)

        self.memory_display = QTextEdit()
        self.memory_display.setReadOnly(True)
        self.memory_display.setStyleSheet("font-size: 10px; background-color: #0F172A; color: #38BDF8;")
        self._update_memory_ui()
        sidebar_layout.addWidget(self.memory_display)

        sidebar_layout.addStretch()
        main_layout.addWidget(self.sidebar)

        self.chat_container = QFrame()
        self.chat_container.setObjectName("ChatContainer")
        chat_layout = QVBoxLayout(self.chat_container)
        chat_layout.setContentsMargins(16, 16, 16, 16)
        chat_layout.setSpacing(12)

        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.viewport().setStyleSheet("background-color: #0F172A;")

        self.chat_scroll_content = QWidget()
        self.chat_scroll_content.setObjectName("ChatScrollContent")
        self.chat_scroll_content.setStyleSheet("background-color: #0F172A;")

        self.chat_scroll_layout = QVBoxLayout(self.chat_scroll_content)
        self.chat_scroll_layout.setAlignment(Qt.AlignmentFlag.AlignTop)
        self.chat_scroll_layout.setSpacing(12)

        self.scroll_area.setWidget(self.chat_scroll_content)
        chat_layout.addWidget(self.scroll_area)

        input_dock = QFrame()
        input_dock.setObjectName("InputDock")
        dock_layout = QVBoxLayout(input_dock)
        dock_layout.setContentsMargins(8, 8, 8, 8)

        self.file_label = QLabel("")
        self.file_label.setStyleSheet("color: #60A5FA; font-size: 11px;")
        dock_layout.addWidget(self.file_label)

        controls_layout = QHBoxLayout()
        self.upload_btn = QPushButton("📎")
        self.upload_btn.setFixedWidth(40)
        self.upload_btn.setObjectName("SecondaryBtn")
        self.upload_btn.clicked.connect(self.select_file_attachment)
        controls_layout.addWidget(self.upload_btn)

        self.user_input = QLineEdit()
        self.user_input.setPlaceholderText("Ask Devgent to build web apps, create presentations with charts, or run commands...")
        self.user_input.returnPressed.connect(self.send_message)
        controls_layout.addWidget(self.user_input)

        self.send_btn = QPushButton("Send")
        self.send_btn.clicked.connect(self.send_message)
        controls_layout.addWidget(self.send_btn)

        dock_layout.addLayout(controls_layout)
        chat_layout.addWidget(input_dock)

        main_layout.addWidget(self.chat_container)

    def _update_memory_ui(self):
        if self.agent_memory:
            self.memory_display.setText(json.dumps(self.agent_memory, indent=2))
        else:
            self.memory_display.setText("Memory is empty.")

    def _on_api_key_changed(self):
        key = self.api_key_input.text().strip()
        if key != self.gemini_api_key:
            self.gemini_api_key = key
            if key:
                save_api_key_to_disk(key)

    def select_output_directory(self):
        chosen = QFileDialog.getExistingDirectory(self, "Select Workspace Folder", self.output_directory)
        if chosen:
            self.output_directory = os.path.normpath(chosen)
            self.path_display.setText(self.output_directory)
            self.memory_file_path = os.path.join(self.output_directory, ".devgent_memory.json")
            self.agent_memory = self._load_memory()
            self._update_memory_ui()

    def select_file_attachment(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Upload", "")
        if file_path:
            self.selected_file_path = os.path.normpath(file_path)
            self.file_label.setText(f"Attached: {os.path.basename(self.selected_file_path)}")

    def _clean_markdown(self, text):
        if not text:
            return ""
        text = text.replace("<br>", "\n").replace("<br/>", "\n")
        text = re.sub(r'```[a-zA-Z]*\n', '', text)
        text = text.replace('```', '')
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        text = re.sub(r'`(.*?)`', r'\1', text)
        text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
        return text.strip()

    def _ensure_extension(self, filename, default_ext):
        filename = str(filename).strip('"\'').strip()
        if not filename:
            return f"output{default_ext}"
        if not os.path.splitext(filename)[1]:
            return f"{filename}{default_ext}"
        return filename

    def _request_confirmation(self, title, message, on_yes, on_no=None):
        """Shows an inline Yes/No permission prompt in the chat window itself (no
        separate popup window) and calls on_yes()/on_no() once the user picks one.
        Used to gate every action that writes a file to disk or runs a command/server."""
        bubble_frame, _ = self.append_chat("Devgent System", f"🔐 {title}\n\n{message}")
        if bubble_frame is None:
            on_yes()
            return

        btn_row = QHBoxLayout()
        btn_row.setContentsMargins(0, 6, 0, 0)
        btn_row.setSpacing(8)

        yes_btn = QPushButton("Yes")
        yes_btn.setObjectName("SuccessBtn")
        no_btn = QPushButton("No")
        no_btn.setObjectName("SecondaryBtn")

        btn_row.addWidget(yes_btn)
        btn_row.addWidget(no_btn)
        btn_row.addStretch()

        bubble_frame.layout().addLayout(btn_row)

        def _resolve(accepted):
            yes_btn.setEnabled(False)
            no_btn.setEnabled(False)
            if accepted:
                on_yes()
            elif on_no:
                on_no()
            QApplication.processEvents()
            self.scroll_area.verticalScrollBar().setValue(self.scroll_area.verticalScrollBar().maximum())

        yes_btn.clicked.connect(lambda: _resolve(True))
        no_btn.clicked.connect(lambda: _resolve(False))

        QApplication.processEvents()
        self.scroll_area.verticalScrollBar().setValue(self.scroll_area.verticalScrollBar().maximum())

    def append_chat(self, sender, message):
        if not message or not str(message).strip():
            return None, None

        bubble_frame = QFrame()
        bubble_layout = QVBoxLayout(bubble_frame)
        bubble_layout.setContentsMargins(12, 10, 12, 10)

        sender_lbl = QLabel(sender)
        sender_lbl.setFont(QFont("Segoe UI", 9, QFont.Weight.Bold))

        if sender == "User":
            bubble_frame.setObjectName("UserBubble")
            bubble_frame.setStyleSheet("QFrame#UserBubble { background-color: #3B82F6; border-radius: 12px; border: none; }")
            sender_lbl.setStyleSheet("color: #E0E7FF; background: transparent; border: none;")
            msg_color = "#FFFFFF"
            align = Qt.AlignmentFlag.AlignRight
        else:
            bubble_frame.setObjectName("DevgentBubble")
            bubble_frame.setStyleSheet("QFrame#DevgentBubble { background-color: #1E293B; border: 1px solid #334155; border-radius: 12px; }")
            sender_lbl.setStyleSheet(("color: #60A5FA;" if sender == "Devgent System" else "color: #94A3B8;") + " background: transparent; border: none;")
            msg_color = "#F8FAFC"
            align = Qt.AlignmentFlag.AlignLeft

        bubble_layout.addWidget(sender_lbl)

        clean_text = self._clean_markdown(message) if sender not in ["User", "Devgent System"] else message
        msg_lbl = QLabel(clean_text)
        msg_lbl.setWordWrap(True)
        msg_lbl.setFont(QFont("Segoe UI", 10))
        msg_lbl.setStyleSheet(f"color: {msg_color}; background: transparent; border: none;")
        bubble_layout.addWidget(msg_lbl)

        msg_key = clean_text[:40]
        self.chat_bubbles[msg_key] = msg_lbl

        wrapper = QHBoxLayout()
        if align == Qt.AlignmentFlag.AlignRight:
            wrapper.addStretch()
            wrapper.addWidget(bubble_frame)
        else:
            wrapper.addWidget(bubble_frame)
            wrapper.addStretch()

        self.chat_scroll_layout.addLayout(wrapper)

        QApplication.processEvents()
        self.scroll_area.verticalScrollBar().setValue(self.scroll_area.verticalScrollBar().maximum())

        return bubble_frame, msg_lbl

    def _add_pptx_footer(self, slide, index, total, accent_color, slide_w, slide_h, on_light_bg=True):
        """Adds a thin accent rule and a 'NN / total' page counter to the bottom of a slide."""
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), slide_h - Inches(0.45), slide_w - Inches(1.0), Pt(1.5))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(accent_color)
        line.line.fill.background()
        line.shadow.inherit = False

        tb = slide.shapes.add_textbox(slide_w - Inches(1.6), slide_h - Inches(0.42), Inches(1.1), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"{index:02d} / {total:02d}"
        p.font.size = Pt(10)
        p.font.name = "Segoe UI"
        p.font.color.rgb = hex_to_rgb(accent_color if on_light_bg else "#94A3B8")
        p.alignment = PP_ALIGN.RIGHT

    def execute_action(self, action_type, action_data):
        print(f"[CONSOLE DEBUG] Executing Action: {action_type}")
        sys.stdout.flush()
        os.makedirs(self.output_directory, exist_ok=True)

        if action_type == "manage_memory":
            # Memory updates are low-risk and local-only, so no permission prompt is needed here.
            updates = action_data.get("updates", {})
            discards = action_data.get("discards", [])

            for k in discards:
                self.agent_memory.pop(k, None)

            for k, v in updates.items():
                self.agent_memory[k] = v

            self._save_memory()
            self._update_memory_ui()
            self.append_chat("Devgent System", "🧠 Agent Memory Updated.")

        elif action_type == "create_batch_files":
            files = action_data.get("files", [])
            names_preview = ", ".join(os.path.basename(self._ensure_extension(f.get("filename", "file.txt"), ".txt")) for f in files)
            def _do_create_batch_files():

                created_count = 0
                for file_info in files:
                    raw_fname = file_info.get("filename", "file.txt")
                    fname = os.path.basename(self._ensure_extension(raw_fname, ".txt"))
                    content = file_info.get("content", "")
                    full_path = os.path.normpath(os.path.join(self.output_directory, fname))
                    try:
                        with open(full_path, "w", encoding="utf-8") as f:
                            f.write(content)
                        created_count += 1
                    except Exception as e:
                        self.append_chat("Devgent System", f"❌ Batch file creation failed for {fname}: {e}")

                if created_count > 0:
                    self.append_chat("Devgent System", f"✅ Created {created_count} file(s) in workspace.")
                    self.start_dev_server()


            self._request_confirmation(
                "Create Files",
                f"Devgent wants to create {len(files)} file(s) in your workspace:\n\n{names_preview}",
                _do_create_batch_files,
                lambda: self.append_chat("Devgent System", "🚫 Batch file creation cancelled by user.")
            )
        elif action_type == "create_txt":
            raw_filename = action_data.get("filename", "output.txt")
            clean_filename = os.path.basename(self._ensure_extension(raw_filename, ".txt"))
            full_save_path = os.path.normpath(os.path.join(self.output_directory, clean_filename))
            def _do_create_txt():
                try:
                    with open(full_save_path, "w", encoding="utf-8") as f:
                        f.write(action_data.get("content", ""))
                    self.append_chat("Devgent System", f"✅ File created successfully:\n{full_save_path}")
                    if clean_filename.endswith((".html", ".htm", ".js", ".css", ".py")):
                        self.start_dev_server()
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Failed to create file: {e}")


            self._request_confirmation(
                "Create File",
                f"Devgent wants to create/overwrite this file:\n\n{clean_filename}",
                _do_create_txt,
                lambda: self.append_chat("Devgent System", "🚫 File creation cancelled by user.")
            )
        elif action_type == "create_docx":
            _ensure_heavy_imports()
            fname = os.path.basename(self._ensure_extension(action_data.get("filename", "Document.docx"), ".docx"))
            full_path = os.path.normpath(os.path.join(self.output_directory, fname))
            def _do_create_docx():
                try:
                    doc = docx.Document()
                    title_text = action_data.get("title", "Document Title")
                    p_title = doc.add_paragraph()
                    run = p_title.add_run(title_text)
                    run.font.name = 'Segoe UI'
                    run.font.size = DocxPt(24)
                    run.font.bold = True
                    run.font.color.rgb = DocxRGBColor(0x1E, 0x3A, 0x8A)

                    # Thin colored accent rule under the title for a bit of visual polish.
                    accent_table = doc.add_table(rows=1, cols=1)
                    accent_cell = accent_table.rows[0].cells[0]
                    accent_cell.width = DocxInches(6.0)
                    shading = OxmlElement('w:shd')
                    shading.set(qn('w:fill'), "3B82F6")
                    accent_cell._tc.get_or_add_tcPr().append(shading)
                    accent_table.rows[0].height = DocxEmu(9525 * 3)

                    if "heading" in action_data:
                        p_head = doc.add_paragraph()
                        r_head = p_head.add_run(action_data["heading"])
                        r_head.font.name = 'Segoe UI'
                        r_head.font.size = DocxPt(16)
                        r_head.font.bold = True
                        r_head.font.color.rgb = DocxRGBColor(0x25, 0x63, 0xEB)

                    for ptext in action_data.get("paragraphs", []):
                        p = doc.add_paragraph()
                        r = p.add_run(ptext)
                        r.font.name = 'Segoe UI'
                        r.font.size = DocxPt(11)

                    doc.save(full_path)
                    self.append_chat("Devgent System", f"📄 Styled Word document created:\n{full_path}")
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Failed to create docx: {e}")


            self._request_confirmation(
                "Create Word Document",
                f"Devgent wants to create a Word document:\n\n{fname}",
                _do_create_docx,
                lambda: self.append_chat("Devgent System", "🚫 Word document creation cancelled by user.")
            )
        elif action_type == "create_pdf":
            _ensure_heavy_imports()
            fname = os.path.basename(self._ensure_extension(action_data.get("filename", "Document.pdf"), ".pdf"))
            full_path = os.path.normpath(os.path.join(self.output_directory, fname))
            def _do_create_pdf():
                try:
                    accent_hex = action_data.get("accent_color", "#3B82F6")
                    accent = colors.HexColor(accent_hex)
                    accent_dark = colors.HexColor(shade_color(accent_hex, -0.35))
                    text_dark = colors.HexColor("#1F2937")
                    muted = colors.HexColor("#64748B")

                    doc = SimpleDocTemplate(
                        full_path, pagesize=letter,
                        topMargin=0.9 * inch, bottomMargin=0.8 * inch,
                        leftMargin=0.75 * inch, rightMargin=0.75 * inch
                    )
                    styles = getSampleStyleSheet()

                    title_style = ParagraphStyle(
                        'CustomTitle', parent=styles['Heading1'],
                        fontSize=26, leading=30, textColor=colors.white,
                        spaceAfter=4, fontName="Helvetica-Bold"
                    )
                    subtitle_style = ParagraphStyle(
                        'CustomSubtitle', parent=styles['Normal'],
                        fontSize=12, leading=16, textColor=colors.HexColor("#E2E8F0")
                    )
                    heading_style = ParagraphStyle(
                        'SectionHeading', parent=styles['Heading2'],
                        fontSize=15, leading=19, textColor=accent_dark,
                        spaceBefore=18, spaceAfter=8, fontName="Helvetica-Bold"
                    )
                    body_style = ParagraphStyle(
                        'CustomBody', parent=styles['Normal'],
                        fontSize=10.5, leading=15, textColor=text_dark, spaceAfter=8
                    )
                    bullet_style = ParagraphStyle(
                        'CustomBullet', parent=body_style, leftIndent=14, spaceAfter=6
                    )
                    quote_style = ParagraphStyle(
                        'CustomQuote', parent=body_style,
                        fontName="Helvetica-Oblique", leftIndent=16, textColor=muted
                    )

                    story = []
                    title_text = action_data.get("title", "Generated Report")
                    subtitle_text = action_data.get("subtitle", "")

                    title_cell_content = [Paragraph(title_text, title_style)]
                    if subtitle_text:
                        title_cell_content.append(Paragraph(subtitle_text, subtitle_style))
                    title_table = Table([[title_cell_content]], colWidths=[6.9 * inch])
                    title_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, -1), accent),
                        ('LEFTPADDING', (0, 0), (-1, -1), 18),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 18),
                        ('TOPPADDING', (0, 0), (-1, -1), 18),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 18),
                    ]))
                    story.append(title_table)
                    story.append(Spacer(1, 18))

                    sections = action_data.get("sections")
                    if not sections:
                        # Backward-compatible: flat paragraphs list.
                        sections = [{"paragraphs": action_data.get("paragraphs", [])}]

                    for section in sections:
                        if section.get("heading"):
                            heading_table = Table([[Paragraph(section["heading"], heading_style)]], colWidths=[6.9 * inch])
                            heading_table.setStyle(TableStyle([
                                ('LINEBELOW', (0, 0), (-1, -1), 1.4, accent),
                                ('LEFTPADDING', (0, 0), (-1, -1), 0),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                            ]))
                            story.append(heading_table)

                        for paragraph in section.get("paragraphs", []):
                            story.append(Paragraph(paragraph, body_style))

                        bullets = section.get("bullets", [])
                        if bullets:
                            items = [ListItem(Paragraph(b, bullet_style), bulletColor=accent, value='●') for b in bullets]
                            story.append(ListFlowable(items, bulletType='bullet', start='●', leftIndent=16))
                            story.append(Spacer(1, 6))

                        if section.get("quote"):
                            q_table = Table([[Paragraph(f"\u201c{section['quote']}\u201d", quote_style)]], colWidths=[6.9 * inch])
                            q_table.setStyle(TableStyle([
                                ('LINEBEFORE', (0, 0), (0, -1), 3, accent),
                                ('LEFTPADDING', (0, 0), (-1, -1), 14),
                                ('TOPPADDING', (0, 0), (-1, -1), 6),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                            ]))
                            story.append(q_table)
                            story.append(Spacer(1, 8))

                        table_data = section.get("table")
                        if table_data and table_data.get("headers") and table_data.get("rows"):
                            t_headers = table_data["headers"]
                            t_rows = table_data["rows"]
                            tbl = Table([t_headers] + t_rows, hAlign='LEFT')
                            tbl.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), accent),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, -1), 9.5),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F5F9")]),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                                ('TOPPADDING', (0, 0), (-1, -1), 6),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                                ('LEFTPADDING', (0, 0), (-1, -1), 8),
                            ]))
                            story.append(tbl)
                            story.append(Spacer(1, 10))

                        story.append(Spacer(1, 4))

                    def _draw_footer(canvas, doc_):
                        canvas.saveState()
                        canvas.setStrokeColor(accent)
                        canvas.setLineWidth(1)
                        canvas.line(0.75 * inch, 0.6 * inch, letter[0] - 0.75 * inch, 0.6 * inch)
                        canvas.setFont("Helvetica", 8)
                        canvas.setFillColor(muted)
                        canvas.drawString(0.75 * inch, 0.45 * inch, title_text)
                        canvas.drawRightString(letter[0] - 0.75 * inch, 0.45 * inch, f"Page {doc_.page}")
                        canvas.restoreState()

                    doc.build(story, onFirstPage=_draw_footer, onLaterPages=_draw_footer)
                    self.append_chat("Devgent System", f"📊 Styled PDF created:\n{full_path}")
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Failed to create PDF: {e}")


            self._request_confirmation(
                "Create PDF",
                f"Devgent wants to create a PDF document:\n\n{fname}",
                _do_create_pdf,
                lambda: self.append_chat("Devgent System", "🚫 PDF creation cancelled by user.")
            )
        elif action_type == "create_pptx":
            _ensure_heavy_imports()
            fname = os.path.basename(self._ensure_extension(action_data.get("filename", "Presentation.pptx"), ".pptx"))
            full_path = os.path.normpath(os.path.join(self.output_directory, fname))
            def _do_create_pptx():
                try:
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

                    theme_accent = action_data.get("theme_accent", action_data.get("title_color", "#3B82F6"))

                    # ---- Title Slide ----
                    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    title_bg_hex = action_data.get("background_color", "#0F172A")
                    bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = hex_to_rgb(title_bg_hex)
                    bg.line.fill.background()
                    bg.shadow.inherit = False

                    # Decorative accent shapes for a less "flat" title slide.
                    accent_circle = title_slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, SLIDE_W - Inches(4.0), Inches(-2.0), Inches(6.0), Inches(6.0)
                    )
                    accent_circle.fill.solid()
                    accent_circle.fill.fore_color.rgb = hex_to_rgb(shade_color(theme_accent, -0.35))
                    accent_circle.line.fill.background()
                    accent_circle.shadow.inherit = False

                    accent_bar = title_slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 0, Inches(2.3), Inches(0.18), Inches(1.9)
                    )
                    accent_bar.fill.solid()
                    accent_bar.fill.fore_color.rgb = hex_to_rgb(theme_accent)
                    accent_bar.line.fill.background()
                    accent_bar.shadow.inherit = False

                    txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(9.5), Inches(3.0))
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    p = tf.paragraphs[0]
                    p.text = action_data.get("title", "Presentation Title")
                    p.font.size = Pt(40)
                    p.font.bold = True
                    p.font.color.rgb = hex_to_rgb(action_data.get("title_color", theme_accent))
                    p.font.name = "Segoe UI"

                    sub = tf.add_paragraph()
                    sub.text = action_data.get("subtitle", "Generated by Devgent Desktop Agent")
                    sub.font.size = Pt(20)
                    sub.font.color.rgb = hex_to_rgb(action_data.get("subtitle_color", "#94A3B8"))
                    sub.font.name = "Segoe UI"
                    sub.space_before = Pt(14)

                    slides_list = action_data.get("slides", [])
                    total_slides = len(slides_list)

                    for i, sdata in enumerate(slides_list, start=1):
                        layout = str(sdata.get("layout", "bullets")).lower()
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        accent_color = sdata.get("accent_color", theme_accent)
                        slide_bg_hex = sdata.get("background_color", "#3B82F6" if layout == "section" else "#F8FAFC")

                        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
                        s_bg.fill.solid()
                        s_bg.fill.fore_color.rgb = hex_to_rgb(slide_bg_hex if layout == "section" else sdata.get("background_color", "#F8FAFC"))
                        s_bg.line.fill.background()
                        s_bg.shadow.inherit = False

                        is_dark_bg = slide_bg_hex.upper() not in ["#FFFFFF", "#F8FAFC", "#F1F5F9", "#FFF"]
                        card_fill = "#1E293B" if is_dark_bg else "#FFFFFF"
                        default_text = sdata.get("text_color", "#F8FAFC" if is_dark_bg else "#334155")

                        if layout == "section":
                            tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.6))
                            tf2 = tb.text_frame
                            tf2.word_wrap = True
                            p2 = tf2.paragraphs[0]
                            p2.text = sdata.get("title", "Section")
                            p2.font.size = Pt(42)
                            p2.font.bold = True
                            p2.font.color.rgb = hex_to_rgb(sdata.get("header_text_color", "#FFFFFF"))
                            p2.font.name = "Segoe UI"
                            p2.alignment = PP_ALIGN.CENTER
                            if sdata.get("subtitle"):
                                sp = tf2.add_paragraph()
                                sp.text = sdata["subtitle"]
                                sp.font.size = Pt(18)
                                sp.font.color.rgb = hex_to_rgb(shade_color(sdata.get("header_text_color", "#FFFFFF"), -0.2))
                                sp.alignment = PP_ALIGN.CENTER
                                sp.font.name = "Segoe UI"
                            self._add_pptx_footer(slide, i, total_slides, "#FFFFFF", SLIDE_W, SLIDE_H, on_light_bg=False)
                            continue

                        # Shared header banner for bullets / two_column / quote layouts.
                        header_bg_hex = sdata.get("header_color", accent_color)
                        header_txt_hex = sdata.get("header_text_color", "#FFFFFF")
                        banner = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), SLIDE_W - Inches(1.0), Inches(0.9)
                        )
                        banner.fill.solid()
                        banner.fill.fore_color.rgb = hex_to_rgb(header_bg_hex)
                        banner.line.fill.background()
                        banner.shadow.inherit = False

                        btf = banner.text_frame
                        btf.word_wrap = True
                        bp = btf.paragraphs[0]
                        bp.text = sdata.get("title", "Slide Title")
                        bp.font.size = Pt(24)
                        bp.font.bold = True
                        bp.font.color.rgb = hex_to_rgb(header_txt_hex)
                        bp.font.name = "Segoe UI"

                        if layout == "quote":
                            quote_card = slide.shapes.add_shape(
                                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.8), SLIDE_W - Inches(2.8), Inches(4.6)
                            )
                            quote_card.fill.solid()
                            quote_card.fill.fore_color.rgb = hex_to_rgb(card_fill)
                            quote_card.line.fill.background()
                            quote_card.shadow.inherit = False

                            q_accent_bar = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(1.8), Inches(0.12), Inches(4.6)
                            )
                            q_accent_bar.fill.solid()
                            q_accent_bar.fill.fore_color.rgb = hex_to_rgb(accent_color)
                            q_accent_bar.line.fill.background()
                            q_accent_bar.shadow.inherit = False

                            qtf = quote_card.text_frame
                            qtf.word_wrap = True
                            qtf.margin_left = Inches(0.5)
                            qtf.margin_right = Inches(0.4)
                            qp = qtf.paragraphs[0]
                            qp.text = f"\u201c{sdata.get('quote_text', '')}\u201d"
                            qp.font.size = Pt(24)
                            qp.font.italic = True
                            qp.font.color.rgb = hex_to_rgb(default_text)
                            qp.font.name = "Segoe UI"
                            if sdata.get("quote_author"):
                                ap = qtf.add_paragraph()
                                ap.text = f"\u2014 {sdata['quote_author']}"
                                ap.font.size = Pt(16)
                                ap.font.bold = True
                                ap.font.color.rgb = hex_to_rgb(accent_color)
                                ap.space_before = Pt(16)
                                ap.font.name = "Segoe UI"

                        elif layout == "two_column":
                            col_w = int((SLIDE_W - Inches(1.4)) / 2)
                            columns = [
                                (Inches(0.5), sdata.get("left_heading", ""), sdata.get("content_left", [])),
                                (Inches(0.5) + col_w + Inches(0.4), sdata.get("right_heading", ""), sdata.get("content_right", [])),
                            ]
                            for left, heading, items in columns:
                                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), col_w, Inches(5.4))
                                card.fill.solid()
                                card.fill.fore_color.rgb = hex_to_rgb(card_fill)
                                card.line.color.rgb = hex_to_rgb(shade_color(card_fill, -0.05))
                                card.shadow.inherit = False

                                ctf = card.text_frame
                                ctf.word_wrap = True
                                ctf.margin_left = Inches(0.3)
                                ctf.margin_top = Inches(0.25)

                                first_para_used = False
                                if heading:
                                    hp = ctf.paragraphs[0]
                                    hp.text = heading
                                    hp.font.bold = True
                                    hp.font.size = Pt(16)
                                    hp.font.color.rgb = hex_to_rgb(accent_color)
                                    hp.font.name = "Segoe UI"
                                    hp.space_after = Pt(10)
                                    first_para_used = True

                                for bullet in items:
                                    bp2 = ctf.add_paragraph() if first_para_used else ctf.paragraphs[0]
                                    first_para_used = True
                                    run1 = bp2.add_run()
                                    run1.text = "\u25cf "
                                    run1.font.color.rgb = hex_to_rgb(accent_color)
                                    run1.font.size = Pt(14)
                                    run1.font.name = "Segoe UI"
                                    run2 = bp2.add_run()
                                    run2.text = bullet
                                    run2.font.color.rgb = hex_to_rgb(default_text)
                                    run2.font.size = Pt(14)
                                    run2.font.name = "Segoe UI"
                                    bp2.space_after = Pt(8)

                        else:  # "bullets" layout (default), optionally with an embedded chart
                            content_items = sdata.get("content", [])
                            has_chart = "chart" in sdata and isinstance(sdata["chart"], dict)

                            if has_chart:
                                text_width = Inches(5.6)
                                chart_left = Inches(6.5)
                                chart_width = SLIDE_W - chart_left - Inches(0.5)
                            else:
                                text_width = SLIDE_W - Inches(1.0)

                            if content_items:
                                card = slide.shapes.add_shape(
                                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), text_width, Inches(5.4)
                                )
                                card.fill.solid()
                                card.fill.fore_color.rgb = hex_to_rgb(card_fill)
                                card.line.color.rgb = hex_to_rgb(shade_color(card_fill, -0.05))
                                card.shadow.inherit = False

                                ctf = card.text_frame
                                ctf.word_wrap = True
                                ctf.margin_left = Inches(0.35)
                                ctf.margin_top = Inches(0.3)
                                for idx, bullet in enumerate(content_items):
                                    cp = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
                                    run1 = cp.add_run()
                                    run1.text = "\u25cf "
                                    run1.font.size = Pt(15)
                                    run1.font.color.rgb = hex_to_rgb(accent_color)
                                    run1.font.name = "Segoe UI"
                                    run2 = cp.add_run()
                                    run2.text = bullet
                                    run2.font.size = Pt(15)
                                    run2.font.color.rgb = hex_to_rgb(default_text)
                                    run2.font.name = "Segoe UI"
                                    cp.space_after = Pt(10)

                            if has_chart:
                                chart_info = sdata["chart"]
                                ctype_str = str(chart_info.get("type", "column")).lower()
                                categories = chart_info.get("categories", [])
                                series_list = chart_info.get("series", [])

                                chart_data = CategoryChartData()
                                chart_data.categories = categories
                                for s in series_list:
                                    chart_data.add_series(s.get("name", "Series"), tuple(s.get("values", [])))

                                chart_type_map = {
                                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                                    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                                    "pie": XL_CHART_TYPE.PIE,
                                    "line": XL_CHART_TYPE.LINE
                                }
                                xl_chart_type = chart_type_map.get(ctype_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

                                chart_shape = slide.shapes.add_chart(
                                    xl_chart_type, chart_left, Inches(1.5), chart_width, Inches(5.4), chart_data
                                )
                                chart = chart_shape.chart
                                chart.has_legend = True
                                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                                chart.legend.include_in_layout = False

                                try:
                                    palette = [accent_color, shade_color(accent_color, 0.3), shade_color(accent_color, -0.3), "#94A3B8"]
                                    plot = chart.plots[0]
                                    for s_idx, series in enumerate(plot.series):
                                        series.format.fill.solid()
                                        series.format.fill.fore_color.rgb = hex_to_rgb(palette[s_idx % len(palette)])
                                except Exception:
                                    pass

                        self._add_pptx_footer(slide, i, total_slides, accent_color, SLIDE_W, SLIDE_H, on_light_bg=not is_dark_bg)

                    prs.save(full_path)
                    self.append_chat("Devgent System", f"💻 Styled PowerPoint created ({total_slides + 1} slides):\n{full_path}")
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Failed to create PPTX: {e}")


            self._request_confirmation(
                "Create PowerPoint",
                f"Devgent wants to create a PowerPoint presentation:\n\n{fname}",
                _do_create_pptx,
                lambda: self.append_chat("Devgent System", "🚫 PowerPoint creation cancelled by user.")
            )
        elif action_type == "create_xlsx":
            _ensure_heavy_imports()
            fname = os.path.basename(self._ensure_extension(action_data.get("filename", "Spreadsheet.xlsx"), ".xlsx"))
            full_path = os.path.normpath(os.path.join(self.output_directory, fname))
            def _do_create_xlsx():
                try:
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = action_data.get("sheet_name", "Sheet1")

                    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
                    header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")

                    alt_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
                    thin_border = Border(
                        left=Side(style='thin', color='CBD5E1'),
                        right=Side(style='thin', color='CBD5E1'),
                        top=Side(style='thin', color='CBD5E1'),
                        bottom=Side(style='thin', color='CBD5E1')
                    )

                    headers = action_data.get("headers", [])
                    if headers:
                        ws.append(headers)
                        for col_num in range(1, len(headers) + 1):
                            cell = ws.cell(row=1, column=col_num)
                            cell.fill = header_fill
                            cell.font = header_font
                            cell.alignment = Alignment(horizontal="center", vertical="center")
                            cell.border = thin_border

                    rows = action_data.get("rows", [])
                    for row_idx, row in enumerate(rows, start=2):
                        ws.append(row)
                        for col_num in range(1, len(row) + 1):
                            cell = ws.cell(row=row_idx, column=col_num)
                            cell.font = Font(name="Segoe UI", size=10)
                            cell.border = thin_border
                            if row_idx % 2 == 0:
                                cell.fill = alt_fill

                    for col in ws.columns:
                        max_len = max(len(str(cell.value or '')) for cell in col)
                        col_letter = get_column_letter(col[0].column)
                        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

                    wb.save(full_path)
                    self.append_chat("Devgent System", f"📈 Styled Excel sheet created ({len(rows)} rows):\n{full_path}")
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Failed to create Excel file: {e}")


            self._request_confirmation(
                "Create Spreadsheet",
                f"Devgent wants to create an Excel spreadsheet:\n\n{fname}",
                _do_create_xlsx,
                lambda: self.append_chat("Devgent System", "🚫 Spreadsheet creation cancelled by user.")
            )
        elif action_type == "run_command":
            cmd = action_data.get("command", "")
            desc = action_data.get("description", "Terminal Command Execution")

            def _do_run_command():
                try:
                    res = subprocess.run(cmd, shell=True, capture_output=True, text=True, cwd=self.output_directory, timeout=60)
                    out = res.stdout.strip() or res.stderr.strip() or "Done with no output."
                    self.append_chat("Devgent System", f"✅ Executed: `{cmd}`\n\nOutput:\n{out}")
                except Exception as e:
                    self.append_chat("Devgent System", f"❌ Command error: {e}")

            self._request_confirmation(
                "Command Permission Required",
                f"Devgent wants to run command:\n\n{desc}\n\n`{cmd}`",
                _do_run_command,
                lambda: self.append_chat("Devgent System", "🚫 Command execution cancelled by user.")
            )

        elif action_type == "start_dev_server":
            self.start_dev_server()

    def start_dev_server(self):
        def _launch():
            os.makedirs(self.output_directory, exist_ok=True)
            index_file = os.path.join(self.output_directory, "index.html")
            if not os.path.exists(index_file):
                with open(index_file, "w", encoding="utf-8") as f:
                    f.write("<!DOCTYPE html><html><body style='background:#0f172a;color:#fff;font-family:sans-serif;'><h1 style='color:#3b82f6;'>Devgent Workspace Online</h1></body></html>")

            if self.active_server_httpd:
                url = f"http://127.0.0.1:{self.active_server_port}"
                webbrowser.open(url)
                self.append_chat("Devgent System", f"🌐 Local Dev Server active at `{url}`")
                return

            ports_to_try = [8000, 8080, 8888, 5000, 3000]
            served_port = None
            httpd = None

            def make_handler(*args, **kwargs):
                return CustomHTTPRequestHandler(*args, directory=self.output_directory, **kwargs)

            for p in ports_to_try:
                try:
                    httpd = ReusableTCPServer(("127.0.0.1", p), make_handler)
                    served_port = p
                    break
                except OSError:
                    continue

            if not httpd:
                self.append_chat("Devgent System", "❌ Failed to bind local server port.")
                return

            self.active_server_httpd = httpd
            self.active_server_port = served_port

            def serve_forever():
                try:
                    httpd.serve_forever()
                except Exception:
                    pass

            threading.Thread(target=serve_forever, daemon=True).start()
            url = f"http://127.0.0.1:{served_port}"
            webbrowser.open(url)
            self.append_chat("Devgent System", f"🌐 Dev Server launched at `{url}`")

        if not self.active_server_httpd:
            self._request_confirmation(
                "Launch Local Server",
                "Devgent wants to start a local web server for this workspace and open it in your browser.",
                _launch,
                lambda: self.append_chat("Devgent System", "🚫 Local server launch cancelled by user.")
            )
            return

        _launch()

    def send_message(self):
        prompt = self.user_input.text().strip()
        attached_file = self.selected_file_path

        if not prompt and not attached_file:
            return

        self.send_btn.setEnabled(False)
        self.user_input.setEnabled(False)

        self.user_input.clear()
        display_text = prompt + (f" [Attached: {os.path.basename(attached_file)}]" if attached_file else "")
        self.append_chat("User", display_text)

        self.chat_history.append({"role": "user", "content": prompt})

        self.selected_file_path = None
        self.file_label.setText("")

        engine_mode = self.engine_option.currentText()
        needs_gemini_key = "gemini-3.6-flash" in engine_mode or "gemini-3.5-flash-lite" in engine_mode

        if needs_gemini_key and not self.gemini_api_key:
            if prompt.strip().lower() == "help":
                self.append_chat("Devgent", GEMINI_API_KEY_HELP_TEXT)
            else:
                self.append_chat("Devgent", "Please enter a free Gemini API key. Reply 'help' if you want to know how to get one.")
            self._unlock_ui()
            return

        self.thinking_bubble_frame, self.thinking_label = self.append_chat("Devgent", "Devgent is thinking... 🧠")

        def async_worker():
            try:
                full_prompt = prompt
                if attached_file:
                    extracted = extract_file_content(attached_file)
                    if extracted:
                        full_prompt += f"\n\n[Attached File Content ({os.path.basename(attached_file)})]:\n" + extracted

                full_prompt += (
                    "\n\n[System reminder: If you learned anything new about the user, "
                    "the project, their preferences, or task status this turn, include a "
                    "manage_memory action to save it - even if the user did not explicitly "
                    "ask you to remember it.]"
                )

                response_text = ""
                sys_inst = self._get_system_instructions()

                if "gemini-3.6-flash" in engine_mode or "gemini-3.5-flash-lite" in engine_mode:
                    if not self.gemini_api_key:
                        return "❌ No Gemini API key configured. Enter a free Gemini API key in the sidebar to add one.", None

                    _ensure_heavy_imports()
                    model_name = "gemini-3.6-flash" if "gemini-3.6-flash" in engine_mode else "gemini-3.5-flash-lite"
                    client = genai.Client(api_key=self.gemini_api_key)

                    config = types.GenerateContentConfig(
                        system_instruction=sys_inst,
                        temperature=0.3
                    )

                    # Convert history for Gemini multi-turn format
                    formatted_contents = []
                    for h in self.chat_history[:-1]:
                        formatted_contents.append(types.Content(
                            role="user" if h["role"] == "user" else "model",
                            parts=[types.Part.from_text(text=h["content"])]
                        ))
                    formatted_contents.append(types.Content(
                        role="user",
                        parts=[types.Part.from_text(text=full_prompt)]
                    ))

                    response = client.models.generate_content(
                        model=model_name,
                        contents=formatted_contents,
                        config=config
                    )
                    response_text = response.text or ""

                elif "Ollama" in engine_mode:
                    msgs = [{"role": "system", "content": sys_inst}]
                    for h in self.chat_history[:-1]:
                        msgs.append({"role": h["role"], "content": h["content"]})
                    msgs.append({"role": "user", "content": full_prompt})

                    resp = requests.post(
                        "http://localhost:11434/v1/chat/completions",
                        json={
                            "model": "deepseek-coder",
                            "messages": msgs
                        },
                        timeout=15
                    )
                    resp.raise_for_status()
                    response_text = resp.json()["choices"][0]["message"]["content"] or ""

                response_text = re.sub(r'<think>[\s\S]*?</think>', '', response_text, flags=re.IGNORECASE).strip()

                payload = None
                raw_json_str = None

                json_match = re.search(r"```json\s*(\{[\s\S]*?\})\s*```", response_text)
                if json_match:
                    raw_json_str = json_match.group(0)
                    extracted_json = json_match.group(1)
                else:
                    fallback_match = re.search(r"(\{[\s\S]*?\})", response_text)
                    if fallback_match:
                        raw_json_str = fallback_match.group(0)
                        extracted_json = fallback_match.group(1)
                    else:
                        extracted_json = None

                if extracted_json:
                    try:
                        payload = json.loads(extracted_json)
                    except Exception:
                        pass

                clean_text = response_text
                if raw_json_str:
                    clean_text = clean_text.replace(raw_json_str, "")

                self.chat_history.append({"role": "assistant", "content": clean_text.strip()})

                return clean_text.strip(), payload

            except Exception as fatal_err:
                return f"❌ Unexpected Failure: {str(fatal_err)}", None

        worker = WorkerThread(async_worker)
        self.active_workers.add(worker)

        def cleanup_worker():
            self.active_workers.discard(worker)
            worker.deleteLater()

        worker.finished_signal.connect(self._on_response_received)
        worker.error_signal.connect(self._on_response_error)
        worker.finished.connect(cleanup_worker)

        worker.start()

    def _unlock_ui(self):
        self.send_btn.setEnabled(True)
        self.user_input.setEnabled(True)
        self.user_input.setFocus()

    def _on_response_received(self, clean_text, payload):
        if self.thinking_label and clean_text:
            self.thinking_label.setText(self._clean_markdown(clean_text))
            self.thinking_label = None
            self.thinking_bubble_frame = None
            QApplication.processEvents()
            self.scroll_area.verticalScrollBar().setValue(self.scroll_area.verticalScrollBar().maximum())
        elif clean_text:
            self.append_chat("Devgent", clean_text)

        if payload:
            actions = payload.get("actions", [])
            if not actions and "action" in payload:
                actions = [payload]

            for act in actions:
                action_type = act.get("action")
                action_data = act.get("action_data", {})
                if action_type:
                    self.execute_action(action_type, action_data)

        self._unlock_ui()

    def _on_response_error(self, err):
        if self.thinking_label:
            self.thinking_label.setText(f"❌ Error: {err}")
            self.thinking_label = None
            self.thinking_bubble_frame = None
        else:
            self.append_chat("Devgent System", f"❌ Error: {err}")
        self._unlock_ui()

    def closeEvent(self, event):
        if self.active_server_httpd:
            try:
                self.active_server_httpd.shutdown()
                self.active_server_httpd.server_close()
            except Exception:
                pass
        event.accept()


if __name__ == "__main__":
    # Kick off the heavy document-generation / Gemini-client imports in a
    # background thread right away, in parallel with the window appearing.
    # By the time the user actually creates a document or sends a message,
    # this has usually already finished.
    threading.Thread(target=_ensure_heavy_imports, daemon=True).start()

    app = QApplication(sys.argv)
    window = DevgentApp()
    window.show()
    sys.exit(app.exec())
